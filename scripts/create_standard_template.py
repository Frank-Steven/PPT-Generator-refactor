"""创建PPT-Generator标准母版模板。

本脚本使用python-pptx创建一个精美的标准母版模板，包含：
- 统一的配色方案（蓝色系商务风格）
- 统一的字体方案
- 页眉页脚设计
- 所有7种标准布局
- 专业的视觉设计
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# 配色方案
PRIMARY_COLOR = RGBColor(0x1F, 0x4E, 0x79)      # 主色：深蓝色
SECONDARY_COLOR = RGBColor(0x2E, 0x75, 0xB6)    # 辅色：中蓝色
ACCENT_COLOR = RGBColor(0x44, 0x72, 0xC4)       # 强调色：亮蓝色
TEXT_DARK = RGBColor(0x33, 0x33, 0x33)          # 深色文本
TEXT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)         # 浅色文本
TEXT_GRAY = RGBColor(0x66, 0x66, 0x66)          # 灰色文本
BG_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)           # 浅色背景
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)           # 白色背景
BORDER_COLOR = RGBColor(0xD0, 0xD7, 0xDE)       # 边框颜色


def set_text_style(
    text_frame,
    font_name: str = "微软雅黑",
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = TEXT_DARK,
    alignment=PP_ALIGN.LEFT,
):
    """设置文本样式。"""
    text_frame.word_wrap = True
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = alignment
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color


def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None):
    """添加形状。"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, **kwargs):
    """添加文本框。"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = text
    set_text_style(tf, **kwargs)
    return txBox


def design_title_slide(slide, prs):
    """设计标题幻灯片布局。"""
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    bg_shape = add_shape(
        slide, MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height,
        fill_color=BG_WHITE
    )

    accent_bar = add_shape(
        slide, MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), slide_height,
        fill_color=PRIMARY_COLOR
    )

    top_bar = add_shape(
        slide, MSO_SHAPE.RECTANGLE, Inches(0.15), 0, slide_width - Inches(0.15), Inches(0.08),
        fill_color=ACCENT_COLOR
    )

    bottom_bar = add_shape(
        slide, MSO_SHAPE.RECTANGLE, Inches(0.15), slide_height - Inches(0.08),
        slide_width - Inches(0.15), Inches(0.08),
        fill_color=SECONDARY_COLOR
    )

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.left = Inches(1.2)
            ph.top = Inches(2.2)
            ph.width = Inches(11)
            ph.height = Inches(1.5)
            for paragraph in ph.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "微软雅黑"
                    run.font.size = Pt(44)
                    run.font.bold = True
                    run.font.color.rgb = PRIMARY_COLOR
        elif ph.placeholder_format.idx == 1:
            ph.left = Inches(1.2)
            ph.top = Inches(3.8)
            ph.width = Inches(10)
            ph.height = Inches(0.8)
            for paragraph in ph.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "微软雅黑"
                    run.font.size = Pt(20)
                    run.font.color.rgb = TEXT_GRAY

    add_text_box(
        slide, Inches(1.2), slide_height - Inches(0.7),
        Inches(8), Inches(0.4),
        "PPT-Generator | 标准母版模板",
        font_size=10, color=TEXT_GRAY
    )


def design_title_and_content_slide(slide, prs):
    """设计标题和内容幻灯片布局。"""
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    header_bg = add_shape(
        slide, MSO_SHAPE.RECTANGLE, 0, 0, slide_width, Inches(1.2),
        fill_color=BG_LIGHT
    )

    accent_line = add_shape(
        slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15),
        Inches(1.5), Inches(0.05),
        fill_color=ACCENT_COLOR
    )

    left_bar = add_shape(
        slide, MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.08), slide_height,
        fill_color=PRIMARY_COLOR
    )

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.left = Inches(0.8)
            ph.top = Inches(0.3)
            ph.width = Inches(12)
            ph.height = Inches(0.7)
            for paragraph in ph.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "微软雅黑"
                    run.font.size = Pt(28)
                    run.font.bold = True
                    run.font.color.rgb = PRIMARY_COLOR
        elif ph.placeholder_format.idx == 1:
            ph.left = Inches(0.8)
            ph.top = Inches(1.5)
            ph.width = Inches(11.5)
            ph.height = Inches(5.5)
            for paragraph in ph.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "微软雅黑"
                    run.font.size = Pt(16)
                    run.font.color.rgb = TEXT_DARK

    add_text_box(
        slide, Inches(0.5), slide_height - Inches(0.5),
        Inches(8), Inches(0.3),
        "PPT-Generator Standard Template",
        font_size=9, color=TEXT_GRAY
    )

    add_text_box(
        slide, slide_width - Inches(1.5), slide_height - Inches(0.5),
        Inches(1), Inches(0.3),
        "",
        font_size=9, color=TEXT_GRAY, alignment=PP_ALIGN.RIGHT
    )


def design_section_header_slide(slide, prs):
    """设计章节标题幻灯片布局。"""
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    bg_shape = add_shape(
        slide, MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height,
        fill_color=BG_LIGHT
    )

    left_panel = add_shape(
        slide, MSO_SHAPE.RECTANGLE, 0, 0, Inches(1.5), slide_height,
        fill_color=PRIMARY_COLOR
    )

    accent_block = add_shape(
        slide, MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3),
        Inches(0.12), Inches(1.2),
        fill_color=ACCENT_COLOR
    )

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.left = Inches(2)
            ph.top = Inches(2.8)
            ph.width = Inches(10)
            ph.height = Inches(1.2)
            for paragraph in ph.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "微软雅黑"
                    run.font.size = Pt(36)
                    run.font.bold = True
                    run.font.color.rgb = PRIMARY_COLOR

    add_text_box(
        slide, Inches(2), Inches(4.2),
        Inches(10), Inches(0.5),
        "Section Header",
        font_size=14, color=TEXT_GRAY
    )


def design_two_content_slide(slide, prs):
    """设计双栏内容幻灯片布局。"""
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    header_bg = add_shape(
        slide, MSO_SHAPE.RECTANGLE, 0, 0, slide_width, Inches(1.2),
        fill_color=BG_LIGHT
    )

    accent_line = add_shape(
        slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15),
        Inches(1.5), Inches(0.05),
        fill_color=ACCENT_COLOR
    )

    left_bar = add_shape(
        slide, MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.08), slide_height,
        fill_color=PRIMARY_COLOR
    )

    divider = add_shape(
        slide, MSO_SHAPE.RECTANGLE, Inches(6.45), Inches(1.5),
        Inches(0.02), Inches(5.3),
        fill_color=BORDER_COLOR
    )

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.left = Inches(0.8)
            ph.top = Inches(0.3)
            ph.width = Inches(12)
            ph.height = Inches(0.7)
            for paragraph in ph.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "微软雅黑"
                    run.font.size = Pt(28)
                    run.font.bold = True
                    run.font.color.rgb = PRIMARY_COLOR
        elif ph.placeholder_format.idx == 1:
            ph.left = Inches(0.8)
            ph.top = Inches(1.5)
            ph.width = Inches(5.4)
            ph.height = Inches(5.3)
            for paragraph in ph.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "微软雅黑"
                    run.font.size = Pt(14)
                    run.font.color.rgb = TEXT_DARK
        elif ph.placeholder_format.idx == 2:
            ph.left = Inches(6.7)
            ph.top = Inches(1.5)
            ph.width = Inches(5.6)
            ph.height = Inches(5.3)
            for paragraph in ph.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "微软雅黑"
                    run.font.size = Pt(14)
                    run.font.color.rgb = TEXT_DARK

    add_text_box(
        slide, Inches(0.5), slide_height - Inches(0.5),
        Inches(8), Inches(0.3),
        "PPT-Generator Standard Template",
        font_size=9, color=TEXT_GRAY
    )


def design_content_with_caption_slide(slide, prs):
    """设计带说明的内容幻灯片布局。"""
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    header_bg = add_shape(
        slide, MSO_SHAPE.RECTANGLE, 0, 0, slide_width, Inches(1.2),
        fill_color=BG_LIGHT
    )

    accent_line = add_shape(
        slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15),
        Inches(1.5), Inches(0.05),
        fill_color=ACCENT_COLOR
    )

    left_bar = add_shape(
        slide, MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.08), slide_height,
        fill_color=PRIMARY_COLOR
    )

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.left = Inches(0.8)
            ph.top = Inches(0.3)
            ph.width = Inches(12)
            ph.height = Inches(0.7)
            for paragraph in ph.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "微软雅黑"
                    run.font.size = Pt(28)
                    run.font.bold = True
                    run.font.color.rgb = PRIMARY_COLOR
        elif ph.placeholder_format.idx == 1:
            ph.left = Inches(0.8)
            ph.top = Inches(1.5)
            ph.width = Inches(11.5)
            ph.height = Inches(4.2)
        elif ph.placeholder_format.idx == 2:
            ph.left = Inches(0.8)
            ph.top = Inches(5.9)
            ph.width = Inches(11.5)
            ph.height = Inches(0.9)
            for paragraph in ph.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "微软雅黑"
                    run.font.size = Pt(12)
                    run.font.color.rgb = TEXT_GRAY

    add_text_box(
        slide, Inches(0.5), slide_height - Inches(0.5),
        Inches(8), Inches(0.3),
        "PPT-Generator Standard Template",
        font_size=9, color=TEXT_GRAY
    )


def design_picture_with_caption_slide(slide, prs):
    """设计带说明的图片幻灯片布局。"""
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    header_bg = add_shape(
        slide, MSO_SHAPE.RECTANGLE, 0, 0, slide_width, Inches(1.2),
        fill_color=BG_LIGHT
    )

    accent_line = add_shape(
        slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15),
        Inches(1.5), Inches(0.05),
        fill_color=ACCENT_COLOR
    )

    left_bar = add_shape(
        slide, MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.08), slide_height,
        fill_color=PRIMARY_COLOR
    )

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.left = Inches(0.8)
            ph.top = Inches(0.3)
            ph.width = Inches(12)
            ph.height = Inches(0.7)
            for paragraph in ph.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "微软雅黑"
                    run.font.size = Pt(28)
                    run.font.bold = True
                    run.font.color.rgb = PRIMARY_COLOR
        elif ph.placeholder_format.idx == 1:
            ph.left = Inches(0.8)
            ph.top = Inches(1.5)
            ph.width = Inches(11.5)
            ph.height = Inches(4.2)
        elif ph.placeholder_format.idx == 2:
            ph.left = Inches(0.8)
            ph.top = Inches(5.9)
            ph.width = Inches(11.5)
            ph.height = Inches(0.9)
            for paragraph in ph.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "微软雅黑"
                    run.font.size = Pt(12)
                    run.font.color.rgb = TEXT_GRAY

    add_text_box(
        slide, Inches(0.5), slide_height - Inches(0.5),
        Inches(8), Inches(0.3),
        "PPT-Generator Standard Template",
        font_size=9, color=TEXT_GRAY
    )


def design_blank_slide(slide, prs):
    """设计空白页幻灯片布局。"""
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    left_bar = add_shape(
        slide, MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.08), slide_height,
        fill_color=PRIMARY_COLOR
    )

    top_bar = add_shape(
        slide, MSO_SHAPE.RECTANGLE, 0, 0, slide_width, Inches(0.04),
        fill_color=ACCENT_COLOR
    )

    add_text_box(
        slide, Inches(0.5), slide_height - Inches(0.5),
        Inches(8), Inches(0.3),
        "PPT-Generator Standard Template",
        font_size=9, color=TEXT_GRAY
    )


def create_standard_template(output_path: Path) -> None:
    """创建标准母版模板。

    参数:
        output_path: 模板文件输出路径。
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    layout_count = len(prs.slide_layouts)

    design_functions = {
        "Title Slide": design_title_slide,
        "Title and Content": design_title_and_content_slide,
        "Section Header": design_section_header_slide,
        "Two Content": design_two_content_slide,
        "Content with Caption": design_content_with_caption_slide,
        "Picture with Caption": design_picture_with_caption_slide,
        "Blank": design_blank_slide,
    }

    layout_info = []
    for i, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for ph in layout.placeholders:
            ph_type = ph.placeholder_format.type
            ph_name = ph.name
            ph_idx = ph.placeholder_format.idx
            placeholders.append(f"{ph_name} (idx={ph_idx}, type={ph_type})")

        layout_info.append({
            "index": i,
            "name": layout.name,
            "placeholders": placeholders,
        })

    for i, layout in enumerate(prs.slide_layouts):
        layout_name = layout.name
        if layout_name in design_functions:
            slide = prs.slides.add_slide(layout)
            design_functions[layout_name](slide, prs)
            print(f"  已设计布局: {layout_name}")
        else:
            print(f"  跳过布局: {layout_name} (无设计函数)")

    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print("\n" + "=" * 60)
    print("标准母版模板创建完成!")
    print("=" * 60)
    print(f"\n文件位置: {output_path}")
    print(f"\n幻灯片尺寸: {prs.slide_width / 914400:.1f} x {prs.slide_height / 914400:.1f} 英寸")
    print(f"\n包含 {layout_count} 个布局:")
    for info in layout_info:
        print(f"\n  [{info['index']}] {info['name']}")
        for ph in info['placeholders']:
            print(f"      - {ph}")


def main() -> None:
    """主函数。"""
    project_root = Path(__file__).parent.parent
    template_path = project_root / "themes" / "standard" / "template.pptx"
    create_standard_template(template_path)


if __name__ == "__main__":
    main()
