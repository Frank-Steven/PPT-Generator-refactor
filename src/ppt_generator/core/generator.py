"""PPT生成器核心模块。

本模块采用函数式工程思想，将生成流程拆分为纯函数管道和IO边界。
所有纯计算逻辑（解析、验证、匹配）与副作用（文件读写、状态变更）彻底分离。

纯函数管道:
    parse_markdown → validate_slides → match_layouts

组合器:
    build_slide_specs 将上述管道与预渲染组合为单一入口

IO边界:
    load_template → render_slides → save_output
"""

from __future__ import annotations

import logging
from collections.abc import Callable
from pathlib import Path

from pptx import Presentation
from returns.pipeline import flow
from returns.result import Failure, Result, Success

from ..matching import LayoutMatcher, PaginationConfig, paginate_slides
from ..parsers import MarkdownParser
from ..prerendering import prerender_slides
from ..rendering import (
    add_slide,
    extract_layouts,
    find_layout_index,
    load_presentation,
    render_image,
    render_list,
    render_paragraph,
    render_rich_list,
    render_rich_paragraph,
    render_slide_item,
    render_title,
    save_presentation,
)
from .exceptions import (
    EmptySlideError,
    MarkdownParseError,
)
from .models import (
    LayoutSpec,
    PrerenderConfig,
    PrerenderResult,
    SlideItem,
    SlideItemType,
    SlideSpec,
    StyleConfig,
    ThemePack,
)

logger = logging.getLogger(__name__)


def render_paragraph_with_style(
    slide: Presentation, item: SlideItem, style_config: StyleConfig
) -> None:
    """渲染段落内容项（带样式）。"""
    runs = item.meta.get("runs", [])
    if runs:
        render_rich_paragraph(slide, runs, style_config)
    else:
        render_paragraph(slide, item.content)


def render_list_with_style(slide: Presentation, item: SlideItem, style_config: StyleConfig) -> None:
    """渲染列表内容项（带样式）。"""
    runs = item.meta.get("runs", [])
    if runs:
        render_rich_list(slide, [runs], style_config)
    else:
        render_list(slide, item.content)


def render_code_block(slide: Presentation, item: SlideItem, style_config: StyleConfig) -> None:
    """渲染代码块内容项。"""
    prerender = item.meta.get("prerender")
    if isinstance(prerender, PrerenderResult):
        render_image(slide, item, style_config)
        return
    render_paragraph(slide, item.content)


def render_table(slide: Presentation, item: SlideItem, style_config: StyleConfig) -> None:
    """渲染表格内容项。"""
    render_paragraph(slide, item.content)


RENDERERS: dict[SlideItemType, Callable[[Presentation, SlideItem, StyleConfig], None]] = {
    SlideItemType.PARAGRAPH: render_paragraph_with_style,
    SlideItemType.LIST: render_list_with_style,
    SlideItemType.CODE: render_code_block,
    SlideItemType.TABLE: render_table,
    SlideItemType.IMAGE: render_image,
}


def parse_markdown(
    markdown_text: str,
    parser: MarkdownParser | None = None,
) -> Result[list[SlideSpec], MarkdownParseError]:
    """解析Markdown文本为幻灯片规格列表。

    参数:
        markdown_text: Markdown源内容。
        parser: 可选的Markdown解析器，默认使用MarkdownParser。

    返回:
        Success(list[SlideSpec])如果解析成功，Failure(MarkdownParseError)否则。
    """
    if not markdown_text.strip():
        return Failure(MarkdownParseError("Markdown源内容不能为空"))

    try:
        parse_fn = parser or MarkdownParser(markdown_text)
        return Success(parse_fn.parse())
    except Exception as exc:
        return Failure(MarkdownParseError(f"解析Markdown失败: {exc}"))


def validate_slides(slides: list[SlideSpec]) -> Result[list[SlideSpec], EmptySlideError]:
    """验证幻灯片列表。

    参数:
        slides: 幻灯片规格列表。

    返回:
        Success(list[SlideSpec])如果验证通过，Failure(EmptySlideError)否则。
    """
    if not slides:
        return Failure(EmptySlideError("未从Markdown中解析到任何幻灯片"))

    for idx, slide in enumerate(slides):
        if not slide.title.strip():
            logger.warning(f"幻灯片 #{idx + 1} 没有标题")

    return Success(slides)


def match_layouts(
    slides: list[SlideSpec],
    layouts: list[LayoutSpec],
    style_config: StyleConfig | None = None,
    matcher: LayoutMatcher | None = None,
) -> list[tuple[SlideSpec, LayoutSpec]]:
    """为每个幻灯片匹配布局。

    参数:
        slides: 幻灯片规格列表。
        layouts: 可用布局列表。
        style_config: 样式配置（用于自定义布局规则）。
        matcher: 布局匹配器，默认使用LayoutMatcher。

    返回:
        (SlideSpec, LayoutSpec)元组列表。
    """
    match_fn = matcher or LayoutMatcher()

    def match_slide(slide: SlideSpec) -> tuple[SlideSpec, LayoutSpec]:
        matched = match_fn.select_layout(slide, layouts)
        default = layouts[0] if layouts else LayoutSpec(name="Default", placeholders=[])
        layout = matched.value_or(default)
        return (slide, layout)

    return [match_slide(slide) for slide in slides]


def build_slide_specs(
    markdown_text: str,
    layouts: list[LayoutSpec],
    style_config: StyleConfig | None = None,
    *,
    parser: MarkdownParser | None = None,
    matcher: LayoutMatcher | None = None,
    prerender_config: PrerenderConfig | None = None,
    pagination_config: PaginationConfig | None = None,
) -> Result[list[tuple[SlideSpec, LayoutSpec]], Exception]:
    """构建完整的幻灯片规格管道。

    将解析、验证、预渲染、匹配、分页组合为单一管道。

    参数:
        markdown_text: Markdown源内容。
        layouts: 可用布局列表。
        style_config: 样式配置。
        parser: 可选的Markdown解析器。
        matcher: 可选的布局匹配器。
        prerender_config: 预渲染配置。
        pagination_config: 分页配置。

    返回:
        Success(list[tuple[SlideSpec, LayoutSpec]])如果成功，Failure(Exception)否则。
    """
    config = style_config or StyleConfig()
    return flow(
        markdown_text,
        lambda text: parse_markdown(text, parser),
        lambda result: result.bind(validate_slides),
        lambda result: result.map(
            lambda slides: _prerender_if_enabled(slides, config, prerender_config)
        ),
        lambda result: result.map(lambda slides: match_layouts(slides, layouts, config, matcher)),
        lambda result: result.map(
            lambda pairs: _paginate_if_enabled(pairs, config, pagination_config)
        ),
    )


def _prerender_if_enabled(
    slides: list[SlideSpec],
    style_config: StyleConfig,
    prerender_config: PrerenderConfig | None,
) -> list[SlideSpec]:
    """根据配置决定是否执行预渲染。

    参数:
        slides: 幻灯片规格列表。
        style_config: 样式配置。
        prerender_config: 预渲染配置。

    返回:
        更新后的幻灯片规格列表。
    """
    if prerender_config:
        logger.info("执行预渲染管线")
        return prerender_slides(slides, style_config, prerender_config)
    return slides


def _paginate_if_enabled(
    slides_with_layouts: list[tuple[SlideSpec, LayoutSpec]],
    style_config: StyleConfig,
    pagination_config: PaginationConfig | None,
) -> list[tuple[SlideSpec, LayoutSpec]]:
    """根据配置决定是否执行自动分页。

    参数:
        slides_with_layouts: (SlideSpec, LayoutSpec) 元组列表。
        style_config: 样式配置。
        pagination_config: 分页配置。

    返回:
        分页后的幻灯片规格列表。
    """
    if pagination_config and pagination_config.enable:
        logger.info("执行自动分页")
        return paginate_slides(slides_with_layouts, style_config, pagination_config)
    return slides_with_layouts


def render_slide(
    presentation: Presentation,
    slide_spec: SlideSpec,
    layout_spec: LayoutSpec,
    style_config: StyleConfig,
    title: str = "Presentation",
) -> None:
    """渲染单个幻灯片。

    参数:
        presentation: Presentation对象。
        slide_spec: 幻灯片规格。
        layout_spec: 布局规格。
        style_config: 样式配置。
        title: 默认标题。
    """
    layout_index = find_layout_index(presentation, layout_spec.name).value_or(0)
    slide = add_slide(presentation, layout_index)

    slide_title = slide_spec.title or title
    render_title(slide, slide_title)

    for item in slide_spec.items:
        render_slide_item(slide, item, style_config, RENDERERS)


def render_presentation(
    presentation: Presentation,
    specs: list[tuple[SlideSpec, LayoutSpec]],
    style_config: StyleConfig,
    title: str = "Presentation",
) -> None:
    """渲染完整演示文稿。

    参数:
        presentation: Presentation对象。
        specs: (SlideSpec, LayoutSpec)元组列表。
        style_config: 样式配置。
        title: 演示文稿标题。
    """
    presentation.core_properties.title = title

    for slide_spec, layout_spec in specs:
        render_slide(presentation, slide_spec, layout_spec, style_config, title)


def generate_ppt(
    markdown_text: str,
    template_path: Path,
    output_path: Path,
    title: str = "Presentation",
    *,
    parser: MarkdownParser | None = None,
    matcher: LayoutMatcher | None = None,
    prerender_config: PrerenderConfig | None = None,
    pagination_config: PaginationConfig | None = None,
) -> Result[Path, Exception]:
    """生成PPT文件的完整流程（使用模板文件）。

    参数:
        markdown_text: Markdown源内容。
        template_path: 模板文件路径。
        output_path: 输出文件路径。
        title: 演示文稿标题。
        parser: 可选的Markdown解析器。
        matcher: 可选的布局匹配器。
        prerender_config: 预渲染配置。
        pagination_config: 分页配置。

    返回:
        Success(Path)如果生成成功，Failure(Exception)否则。
    """
    return _generate(
        markdown_text=markdown_text,
        template_path=template_path,
        output_path=output_path,
        title=title,
        style_config=StyleConfig(),
        parser=parser,
        matcher=matcher,
        prerender_config=prerender_config,
        pagination_config=pagination_config,
    )


def generate_ppt_with_theme(
    markdown_text: str,
    theme_pack: ThemePack,
    output_path: Path,
    title: str = "Presentation",
    *,
    parser: MarkdownParser | None = None,
    matcher: LayoutMatcher | None = None,
    prerender_config: PrerenderConfig | None = None,
    pagination_config: PaginationConfig | None = None,
) -> Result[Path, Exception]:
    """生成PPT文件的完整流程（使用主题包）。

    参数:
        markdown_text: Markdown源内容。
        theme_pack: 主题包。
        output_path: 输出文件路径。
        title: 演示文稿标题。
        parser: 可选的Markdown解析器。
        matcher: 可选的布局匹配器。
        prerender_config: 预渲染配置。
        pagination_config: 分页配置。

    返回:
        Success(Path)如果生成成功，Failure(Exception)否则。
    """
    effective_matcher = matcher or LayoutMatcher(layout_config=theme_pack.layout_config)
    return _generate(
        markdown_text=markdown_text,
        template_path=theme_pack.template_path,
        output_path=output_path,
        title=title,
        style_config=theme_pack.style_config,
        parser=parser,
        matcher=effective_matcher,
        prerender_config=prerender_config,
        pagination_config=pagination_config,
    )


def _generate(
    markdown_text: str,
    template_path: Path,
    output_path: Path,
    title: str,
    style_config: StyleConfig,
    parser: MarkdownParser | None,
    matcher: LayoutMatcher | None,
    prerender_config: PrerenderConfig | None,
    pagination_config: PaginationConfig | None = None,
) -> Result[Path, Exception]:
    """生成PPT的内部实现，统一处理模板和主题包两种模式。"""
    return flow(
        template_path,
        extract_layouts,
        lambda result: result.bind(
            lambda layouts: _build_and_render(
                layouts,
                markdown_text,
                template_path,
                output_path,
                title,
                style_config,
                parser,
                matcher,
                prerender_config,
                pagination_config,
            )
        ),
    )


def _build_and_render(
    layouts: list[LayoutSpec],
    markdown_text: str,
    template_path: Path,
    output_path: Path,
    title: str,
    style_config: StyleConfig,
    parser: MarkdownParser | None,
    matcher: LayoutMatcher | None,
    prerender_config: PrerenderConfig | None,
    pagination_config: PaginationConfig | None = None,
) -> Result[Path, Exception]:
    """构建幻灯片规格并渲染输出。

    使用 Result.bind 组合管道，消除 isinstance(Failure) 检查。
    所有错误通过 bind 自动短路传播。

    参数:
        layouts: 布局规格列表。
        markdown_text: Markdown源内容。
        template_path: 模板文件路径。
        output_path: 输出文件路径。
        title: 演示文稿标题。
        style_config: 样式配置。
        parser: Markdown解析器。
        matcher: 布局匹配器。
        prerender_config: 预渲染配置。
        pagination_config: 分页配置。

    返回:
        Success(Path)如果生成成功，Failure(Exception)否则。
    """
    return build_slide_specs(
        markdown_text,
        layouts,
        style_config,
        parser=parser,
        matcher=matcher,
        prerender_config=prerender_config,
        pagination_config=pagination_config,
    ).bind(lambda specs: _render_and_save(specs, template_path, output_path, style_config, title))


def _render_and_save(
    specs: list[tuple[SlideSpec, LayoutSpec]],
    template_path: Path,
    output_path: Path,
    style_config: StyleConfig,
    title: str,
) -> Result[Path, Exception]:
    """渲染幻灯片并保存文件。

    通过 bind 链式组合 load → render → save，错误自动短路。

    参数:
        specs: (SlideSpec, LayoutSpec)元组列表。
        template_path: 模板文件路径。
        output_path: 输出文件路径。
        style_config: 样式配置。
        title: 演示文稿标题。

    返回:
        Success(Path)如果保存成功，Failure(Exception)否则。
    """
    return (
        load_presentation(template_path)
        .bind(
            lambda presentation: _try_render(presentation, specs, style_config, title).map(
                lambda _: presentation
            )
        )
        .bind(
            lambda presentation: save_presentation(presentation, output_path).map(
                lambda _: _log_success(output_path)
            )
        )
    )


def _log_success(output_path: Path) -> Path:
    """记录成功日志并返回路径，供 Result.map 链式调用使用。"""
    logger.info(f"演示文稿已成功生成: {output_path}")
    return output_path


def _try_render(
    presentation: Presentation,
    specs: list[tuple[SlideSpec, LayoutSpec]],
    style_config: StyleConfig,
    title: str,
) -> Result[None, Exception]:
    """执行渲染操作，将可能的异常包装为 Result。

    参数:
        presentation: Presentation对象。
        specs: 幻灯片规格列表。
        style_config: 样式配置。
        title: 演示文稿标题。

    返回:
        Success(None)如果渲染成功，Failure(Exception)否则。
    """
    try:
        render_presentation(presentation, specs, style_config, title)
        return Success(None)
    except Exception as exc:
        return Failure(exc)


class PPTGenerator:
    """从Markdown和PPT模板生成PowerPoint演示文稿的主类。

    保持向后兼容性，内部使用函数式管道实现。支持模板文件和主题包两种模式。
    """

    def __init__(
        self,
        markdown_text: str,
        template_path: Path | None = None,
        output_path: Path | None = None,
        title: str = "Presentation",
        *,
        theme_pack: ThemePack | None = None,
        parser: MarkdownParser | None = None,
        matcher: LayoutMatcher | None = None,
        prerender_config: PrerenderConfig | None = None,
        pagination_config: PaginationConfig | None = None,
    ) -> None:
        self.markdown_text = markdown_text
        self.template_path = Path(template_path) if template_path else None
        self.output_path = Path(output_path) if output_path else None
        self.title = title
        self._theme_pack = theme_pack
        self._parser = parser
        self._matcher = matcher
        self._prerender_config = prerender_config
        self._pagination_config = pagination_config
        self._result: Result[Path, Exception] | None = None

    def generate(self) -> None:
        """执行完整的生成管道。

        Raises:
            PPTGeneratorError: 如果生成过程中发生错误。
            ValueError: 如果未提供 template_path 或 theme_pack。
        """
        if self._theme_pack:
            self._result = generate_ppt_with_theme(
                self.markdown_text,
                self._theme_pack,
                self.output_path or Path("output.pptx"),
                self.title,
                parser=self._parser,
                matcher=self._matcher,
                prerender_config=self._prerender_config,
                pagination_config=self._pagination_config,
            )
        elif self.template_path:
            self._result = generate_ppt(
                self.markdown_text,
                self.template_path,
                self.output_path or Path("output.pptx"),
                self.title,
                parser=self._parser,
                matcher=self._matcher,
                prerender_config=self._prerender_config,
                pagination_config=self._pagination_config,
            )
        else:
            raise ValueError("必须提供 template_path 或 theme_pack")

        if isinstance(self._result, Failure):
            raise self._result.failure()

    @property
    def result(self) -> Result[Path, Exception] | None:
        """返回生成结果。"""
        return self._result

    def metadata_json(self) -> bytes:
        """将演示文稿元数据导出为JSON。"""
        import orjson

        data: bytes = orjson.dumps(
            {
                "title": self.title,
            }
        )
        return data
