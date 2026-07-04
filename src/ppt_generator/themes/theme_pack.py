"""主题包加载器和验证工具。

本模块提供主题包的加载、验证和解析功能，支持从主题包目录加载完整的主题定义。
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from ..core.exceptions import InvalidConfigError, TemplateLoadError, MissingFileError
from ..core.models import (
    ThemePack,
    ThemePackManifest,
    StyleConfig,
    CodeStyle,
    MermaidStyle,
    LatexStyle,
    TableStyle,
    RunStyle,
    RunOverrides,
    LayoutConfig,
    LayoutDef,
    LayoutGroupDef,
    LayoutDefaults,
    LayoutPlaceholderDef,
    LayoutAutoApply,
)
from ..utils import load_yaml_file
from pptx import Presentation

logger = logging.getLogger(__name__)


def load_theme_pack(theme_pack_path: Path | str) -> ThemePack:
    """加载主题包。

    Args:
        theme_pack_path: 主题包目录路径。

    Returns:
        ThemePack实例。

    Raises:
        MissingFileError: 如果主题包目录或必需文件不存在。
        InvalidConfigError: 如果主题包配置无效。
        TemplateLoadError: 如果模板验证失败。
    """
    path = Path(theme_pack_path)
    
    if not path.exists():
        raise MissingFileError(f"主题包目录不存在: {path}")
    
    if not path.is_dir():
        raise InvalidConfigError(f"主题包路径必须是目录: {path}")
    
    manifest = _load_manifest(path)
    template_path = _resolve_template_path(path, manifest)
    style_config = _load_style_config(path, manifest)
    layout_config = _load_layout_config(path, manifest)
    
    _validate_template(template_path, layout_config)
    
    preview_path = _resolve_optional_path(path, manifest.files.get("preview"))
    fonts_path = path / "fonts" if (path / "fonts").exists() else None
    assets_path = path / "assets" if (path / "assets").exists() else None
    
    return ThemePack(
        manifest=manifest,
        template_path=template_path,
        style_config=style_config,
        layout_config=layout_config,
        preview_path=preview_path,
        fonts_path=fonts_path,
        assets_path=assets_path,
    )


def _load_manifest(theme_pack_path: Path) -> ThemePackManifest:
    """加载manifest.yaml文件。"""
    manifest_path = theme_pack_path / "manifest.yaml"
    data = load_yaml_file(manifest_path, "manifest.yaml")
    return ThemePackManifest(
        name=data.get("name", ""),
        version=data.get("version", ""),
        author=data.get("author", ""),
        description=data.get("description"),
        spec_version=data.get("spec_version", "1.0"),
        compatible_generator=data.get("compatible_generator"),
        files=data.get("files", {}),
        preview=data.get("preview", {}),
        tags=data.get("tags", []),
    )


def _resolve_template_path(theme_pack_path: Path, manifest: ThemePackManifest) -> Path:
    """解析模板文件路径。"""
    template_filename = manifest.files.get("template", "template.pptx")
    template_path = theme_pack_path / template_filename
    
    if not template_path.exists():
        raise MissingFileError(f"模板文件未找到: {template_path}")
    
    return template_path


def _resolve_optional_path(theme_pack_path: Path, filename: str | None) -> Path | None:
    """解析可选文件路径。"""
    if not filename:
        return None
    
    path = theme_pack_path / filename
    return path if path.exists() else None


def _load_style_config(theme_pack_path: Path, manifest: ThemePackManifest) -> StyleConfig:
    """加载style.yaml配置文件。"""
    style_filename = manifest.files.get("style", "style.yaml")
    style_path = theme_pack_path / style_filename
    
    if not style_path.exists():
        return StyleConfig()
    
    data = load_yaml_file(style_path, "style.yaml")
    return _parse_style_config(data)


def _parse_style_config(data: dict[str, Any]) -> StyleConfig:
    """解析样式配置数据。"""
    code_data = data.get("code", {})
    code_style = CodeStyle(
        font=code_data.get("font", "Consolas"),
        font_size=code_data.get("font_size", 11),
        theme=code_data.get("theme", "monokai"),
        line_numbers=code_data.get("line_numbers", True),
        background_color=code_data.get("background_color", "#272822"),
        text_color=code_data.get("text_color", "#F8F8F2"),
        border_radius=code_data.get("border_radius", 4),
        padding=code_data.get("padding", 12),
        line_height=code_data.get("line_height", 1.4),
    )
    
    mermaid_data = data.get("mermaid", {})
    mermaid_style = MermaidStyle(
        theme=mermaid_data.get("theme", "dark"),
        background_color=mermaid_data.get("background_color", "#1a1a1a"),
        scale=mermaid_data.get("scale", 2),
        padding=mermaid_data.get("padding", 10),
    )
    
    latex_data = data.get("latex", {})
    latex_style = LatexStyle(
        font_size=latex_data.get("font_size", 14),
        background_color=latex_data.get("background_color", "transparent"),
        dpi=latex_data.get("dpi", 300),
        color=latex_data.get("color", "#333333"),
    )
    
    table_data = data.get("table", {})
    table_style = TableStyle(
        font=table_data.get("font", "微软雅黑"),
        font_size=table_data.get("font_size", 10),
        header_background=table_data.get("header_background", "#4472C4"),
        header_color=table_data.get("header_color", "#FFFFFF"),
        even_row_background=table_data.get("even_row_background", "#F5F5F5"),
        odd_row_background=table_data.get("odd_row_background", "#FFFFFF"),
        border_color=table_data.get("border_color", "#CCCCCC"),
        border_width=table_data.get("border_width", 1),
    )
    
    run_overrides_data = data.get("run_overrides", {})
    run_overrides = _parse_run_overrides(run_overrides_data)
    
    return StyleConfig(
        code=code_style,
        mermaid=mermaid_style,
        latex=latex_style,
        table=table_style,
        run_overrides=run_overrides,
    )


def _load_layout_config(theme_pack_path: Path, manifest: ThemePackManifest) -> LayoutConfig:
    """加载layouts.yaml配置文件。"""
    layouts_filename = manifest.files.get("layouts", "layouts.yaml")
    layouts_path = theme_pack_path / layouts_filename
    
    if not layouts_path.exists():
        raise MissingFileError(f"layouts.yaml未找到: {layouts_path}")
    
    data = load_yaml_file(layouts_path, "layouts.yaml")
    return _parse_layout_config(data)


def _parse_layout_config(data: dict[str, Any]) -> LayoutConfig:
    """解析布局配置数据。"""
    version = data.get("version", "1.0")
    
    defaults_data = data.get("defaults", {})
    defaults = LayoutDefaults(
        default=defaults_data.get("default", "title-and-content"),
        first_slide=defaults_data.get("first_slide", "title-slide"),
        section_divider=defaults_data.get("section_divider", "section-header"),
        content=defaults_data.get("content", "title-and-content"),
        multi_column=defaults_data.get("multi_column", "two-content"),
        media=defaults_data.get("media", "content-with-caption"),
        image=defaults_data.get("image", "picture-with-caption"),
        full_width=defaults_data.get("full_width", "blank"),
    )
    
    groups_data = data.get("groups", {})
    groups = {}
    for group_id, group_data in groups_data.items():
        if not isinstance(group_data, dict):
            continue
        groups[group_id] = LayoutGroupDef(
            id=group_id,
            display_name=group_data.get("display_name"),
            description=group_data.get("description"),
        )
    
    layouts_data = data.get("layouts", [])
    layouts = []
    for layout_data in layouts_data:
        if not isinstance(layout_data, dict):
            continue
        layouts.append(_parse_layout_def(layout_data))
    
    return LayoutConfig(
        version=version,
        defaults=defaults,
        groups=groups,
        layouts=layouts,
    )


def _parse_layout_def(data: dict[str, Any]) -> LayoutDef:
    """解析单个布局定义。"""
    placeholders_data = data.get("placeholders", [])
    placeholders = []
    for ph_data in placeholders_data:
        if not isinstance(ph_data, dict):
            continue
        placeholders.append(LayoutPlaceholderDef(
            index=ph_data.get("index", 0),
            type=ph_data.get("type", ""),
            role=ph_data.get("role"),
            name=ph_data.get("name"),
            description=ph_data.get("description"),
        ))
    
    auto_apply_data = data.get("auto_apply")
    auto_apply = None
    if isinstance(auto_apply_data, dict):
        auto_apply = LayoutAutoApply(
            conditions=auto_apply_data.get("conditions", []),
            priority=auto_apply_data.get("priority", 0),
        )
    elif isinstance(auto_apply_data, list):
        auto_apply = LayoutAutoApply(conditions=auto_apply_data)
    
    return LayoutDef(
        id=data.get("id", ""),
        name=data.get("name", ""),
        display_name=data.get("display_name"),
        description=data.get("description"),
        group=data.get("group"),
        placeholders=placeholders,
        keywords=data.get("keywords", []),
        tags=data.get("tags", []),
        auto_apply=auto_apply,
    )


def _parse_run_overrides(data: dict[str, Any]) -> RunOverrides:
    """解析Run级别样式覆盖。"""
    bold_style = _parse_run_style(data.get("bold", {}))
    italic_style = _parse_run_style(data.get("italic", {}))
    code_style = _parse_run_style(data.get("code", {}))
    link_style = _parse_run_style(data.get("link", {}))
    
    return RunOverrides(
        bold=bold_style,
        italic=italic_style,
        code=code_style,
        link=link_style,
    )


def _parse_run_style(data: dict[str, Any]) -> RunStyle:
    """解析单个Run样式。"""
    return RunStyle(
        font=data.get("font"),
        font_size=data.get("font_size"),
        color=data.get("color"),
        bold=data.get("bold"),
        italic=data.get("italic"),
        underline=data.get("underline"),
        background_color=data.get("background_color"),
    )


def _validate_template(template_path: Path, layout_config: LayoutConfig) -> None:
    """验证模板是否与布局配置一致。

    检查模板中是否包含layouts.yaml中定义的所有布局。
    """
    try:
        prs = Presentation(str(template_path))
    except Exception as e:
        raise TemplateLoadError(f"无法加载模板文件: {e}") from e
    
    template_layout_names = {layout.name for layout in prs.slide_layouts}
    
    defined_layout_names = []
    for layout_def in layout_config.layouts:
        defined_layout_names.append(layout_def.name)
        if layout_def.name not in template_layout_names:
            raise InvalidConfigError(
                f"模板缺少layouts.yaml中定义的布局: '{layout_def.name}' "
                f"(id: {layout_def.id})"
            )


def list_available_themes(themes_dir: Path | str) -> list[ThemePackManifest]:
    """列出指定目录下所有可用的主题包。

    Args:
        themes_dir: 主题包存放目录。

    Returns:
        主题包元数据列表。
    """
    path = Path(themes_dir)
    
    if not path.exists():
        return []
    
    manifests = []
    for item in path.iterdir():
        if item.is_dir():
            manifest_path = item / "manifest.yaml"
            if manifest_path.exists():
                try:
                    manifest = _load_manifest(item)
                    manifests.append(manifest)
                except Exception as exc:
                    logger.warning(f"加载主题包失败，跳过: {item}", exc_info=True)
                    continue
    
    return manifests
