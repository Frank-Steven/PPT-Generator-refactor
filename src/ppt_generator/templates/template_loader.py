"""PPT模板加载器。

本模块负责加载和验证PPT模板文件。
"""

from __future__ import annotations

from pathlib import Path
from typing import Optional

from pptx import Presentation

from ..core.exceptions import TemplateLoadError, MissingFileError
from ..core.models import LayoutSpec, PlaceholderSpec
from ..utils import placeholder_type_to_str


class TemplateLoader:
    """PPT模板加载器。

    负责加载PPT模板文件，并提供布局信息查询功能。
    """

    def __init__(self, template_path: Path | str) -> None:
        """初始化模板加载器。

        参数:
            template_path: PPT模板文件路径。

        异常:
            MissingFileError: 如果模板文件不存在。
            TemplateLoadError: 如果模板无法加载。
        """
        self._path = Path(template_path)
        
        if not self._path.exists():
            raise MissingFileError(f"模板文件不存在: {self._path}")
        
        try:
            self._presentation = Presentation(str(self._path))
        except Exception as exc:
            raise TemplateLoadError(f"加载模板失败: {exc}") from exc

    @property
    def presentation(self) -> Presentation:
        """获取Presentation对象。

        返回:
            Presentation对象。
        """
        return self._presentation

    @property
    def path(self) -> Path:
        """获取模板文件路径。

        返回:
            模板文件路径。
        """
        return self._path

    def list_layouts(self) -> list[LayoutSpec]:
        """列出模板中所有可用的布局。

        返回:
            布局规格列表。
        """
        layouts = []
        for layout in self._presentation.slide_layouts:
            placeholders = []
            for ph in layout.placeholders:
                placeholders.append(PlaceholderSpec(
                    name=ph.name,
                    placeholder_type=placeholder_type_to_str(ph.placeholder_format.type),
                    index=ph.placeholder_format.idx,
                    shape_id=ph.shape_id,
                ))
            layouts.append(LayoutSpec(
                name=layout.name,
                placeholders=placeholders,
            ))
        return layouts

    def get_layout_by_name(self, name: str) -> Optional[LayoutSpec]:
        """根据名称获取布局。

        参数:
            name: 布局名称。

        返回:
            布局规格，如果找不到则返回None。
        """
        for layout in self.list_layouts():
            if layout.name == name:
                return layout
        return None

    def validate(self) -> bool:
        """验证模板是否符合标准规范。

        返回:
            如果模板有效返回True，否则返回False。
        """
        required_layouts = [
            "Title Slide",
            "Title and Content",
            "Section Header",
        ]
        
        layout_names = [layout.name for layout in self.list_layouts()]
        
        for required in required_layouts:
            if required not in layout_names:
                return False
        
        return True