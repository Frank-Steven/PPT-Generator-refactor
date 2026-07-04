"""IO副作用模块。

本模块收敛所有与python-pptx相关的可变操作，实现IO边界分离。
所有函数都是纯副作用函数，不返回有意义的值或返回Result类型。

设计原则:
1. 所有对外部世界的修改（文件读写、状态变更）都放在这里
2. 纯计算逻辑在其他模块中，通过数据传递与本模块交互
3. 使用Result类型处理IO错误，避免抛出异常
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Callable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.slide import Slide
from pptx.util import Pt
from returns.result import Failure, Result, Success

from ..core.exceptions import TemplateLoadError
from ..core.models import (
    SlideItem,
    SlideItemType,
    LayoutSpec,
    RichRun,
    StyleConfig,
    RunStyle,
    PrerenderResult,
)
from PIL import Image as PILImage

from ..templates import TemplateLoader
from ..utils import ensure_dir, hex_to_rgb


def load_presentation(template_path: Path) -> Result[Presentation, TemplateLoadError]:
    """加载PPT模板文件。

    参数:
        template_path: 模板文件路径。

    返回:
        Success(Presentation)如果加载成功，Failure(TemplateLoadError)否则。
    """
    try:
        loader = TemplateLoader(template_path)
        return Success(loader.presentation)
    except TemplateLoadError as exc:
        return Failure(exc)
    except Exception as exc:
        return Failure(TemplateLoadError(f"加载模板失败: {exc}"))


def save_presentation(presentation: Presentation, output_path: Path) -> Result[None, Exception]:
    """保存演示文稿到输出路径。

    参数:
        presentation: Presentation对象。
        output_path: 输出文件路径。

    返回:
        Success(None)如果保存成功，Failure(Exception)否则。
    """
    try:
        ensure_dir(output_path.parent)
        presentation.save(str(output_path))
        return Success(None)
    except Exception as exc:
        return Failure(exc)


def find_layout_index(presentation: Presentation, layout_name: str) -> Result[int, ValueError]:
    """根据布局名称查找索引。

    参数:
        presentation: Presentation对象。
        layout_name: 布局名称。

    返回:
        Success(int)如果找到，Failure(ValueError)否则。
    """
    index = next(
        (i for i, layout in enumerate(presentation.slide_layouts) if layout.name == layout_name),
        None,
    )
    if index is not None:
        return Success(index)
    return Failure(ValueError(f"布局 '{layout_name}' 未找到"))


def add_slide(presentation: Presentation, layout_index: int) -> Slide:
    """添加新幻灯片。

    参数:
        presentation: Presentation对象。
        layout_index: 布局索引。

    返回:
        新创建的Slide对象。
    """
    return presentation.slides.add_slide(presentation.slide_layouts[layout_index])


def render_title(slide: Slide, title: str) -> None:
    """渲染幻灯片标题。

    优先匹配 TITLE 和 CENTER_TITLE 类型的占位符，避免匹配 SUBTITLE。

    参数:
        slide: Slide对象。
        title: 标题文本。
    """
    title_types = (PP_PLACEHOLDER_TYPE.TITLE, PP_PLACEHOLDER_TYPE.CENTER_TITLE)

    shape = next(
        (s for s in slide.shapes
         if s.is_placeholder and s.placeholder_format.type in title_types),
        None,
    )
    if shape is not None:
        shape.text = title
        return

    shape = next(
        (s for s in slide.shapes
         if s.is_placeholder
         and "title" in s.name.lower()
         and s.placeholder_format.type != PP_PLACEHOLDER_TYPE.SUBTITLE),
        None,
    )
    if shape is not None:
        shape.text = title


def _find_first_placeholder(slide: Slide, ph_types: tuple[PP_PLACEHOLDER_TYPE, ...]) -> Any | None:
    """按类型优先级查找第一个匹配的占位符文本框。

    高阶辅助函数，消除 get_body_text_frame 中的三重 for 循环。

    参数:
        slide: Slide对象。
        ph_types: 占位符类型元组，按优先级排序。

    返回:
        TextFrame对象，如果找不到则返回None。
    """
    for ph_type in ph_types:
        shape = next(
            (s for s in slide.shapes
             if s.is_placeholder
             and s.placeholder_format.type == ph_type
             and s.has_text_frame),
            None,
        )
        if shape is not None:
            return shape.text_frame
    return None


def get_body_text_frame(slide: Slide) -> Any | None:
    """获取幻灯片的正文文本框。

    按优先级查找：BODY > OBJECT > SUBTITLE。

    参数:
        slide: Slide对象。

    返回:
        TextFrame对象，如果找不到则返回None。
    """
    return _find_first_placeholder(slide, (
        PP_PLACEHOLDER_TYPE.BODY,
        PP_PLACEHOLDER_TYPE.OBJECT,
        PP_PLACEHOLDER_TYPE.SUBTITLE,
    ))


def set_autofit(text_frame: Any) -> None:
    """设置文本框自动调整大小。

    参数:
        text_frame: TextFrame对象。
    """
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def apply_run_style(run: Any, run_style: RunStyle) -> None:
    """应用Run级别样式覆盖。

    参数:
        run: python-pptx Run对象。
        run_style: Run样式配置。
    """
    if run_style.font is not None:
        run.font.name = run_style.font
    if run_style.font_size is not None:
        run.font.size = Pt(run_style.font_size)
    if run_style.color is not None:
        run.font.color.rgb = _hex_to_rgb(run_style.color)
    if run_style.bold is not None:
        run.bold = run_style.bold
    if run_style.italic is not None:
        run.italic = run_style.italic
    if run_style.underline is not None:
        run.font.underline = run_style.underline
    if run_style.background_color is not None:
        run.font.highlight_color = _hex_to_rgb(run_style.background_color)


def render_rich_paragraph(
    slide: Slide,
    runs: list[RichRun],
    style_config: StyleConfig,
    append: bool = True,
) -> None:
    """渲染带样式的段落内容。

    参数:
        slide: Slide对象。
        runs: RichRun列表。
        style_config: 样式配置。
        append: 是否追加到现有内容，默认为True。
    """
    text_frame = get_body_text_frame(slide)
    if text_frame is None:
        return
    
    set_autofit(text_frame)
    
    if not append or not text_frame.paragraphs:
        text_frame.clear()
    
    p = text_frame.add_paragraph()
    
    for run_data in runs:
        run = p.add_run()
        run.text = run_data.text
        
        if run_data.bold:
            run.bold = True
            apply_run_style(run, style_config.run_overrides.bold)
        
        if run_data.italic:
            run.italic = True
            apply_run_style(run, style_config.run_overrides.italic)
        
        if run_data.code:
            apply_run_style(run, style_config.run_overrides.code)
        
        if run_data.strikethrough:
            run.font.strikethrough = True
        
        if run_data.link:
            apply_run_style(run, style_config.run_overrides.link)


def render_paragraph(slide: Slide, content: str, append: bool = True) -> None:
    """渲染段落内容。

    参数:
        slide: Slide对象。
        content: 段落文本。
        append: 是否追加到现有内容，默认为True。
    """
    text_frame = get_body_text_frame(slide)
    if text_frame is None:
        return
    
    set_autofit(text_frame)
    
    if not append or not text_frame.paragraphs:
        text_frame.clear()
    
    p = text_frame.add_paragraph()
    p.text = content


def render_list(slide: Slide, content: str, level: int = 1) -> None:
    """渲染列表内容。

    参数:
        slide: Slide对象。
        content: 列表项内容。
        level: 列表级别（0为顶级）。
    """
    text_frame = get_body_text_frame(slide)
    if text_frame is None:
        return
    
    set_autofit(text_frame)
    
    if not text_frame.paragraphs:
        text_frame.clear()
    
    p = text_frame.add_paragraph()
    p.text = content
    p.level = level


def render_rich_list(
    slide: Slide,
    items: list[list[RichRun]],
    style_config: StyleConfig,
    level: int = 1,
) -> None:
    """渲染带样式的列表内容。

    参数:
        slide: Slide对象。
        items: 每个列表项的RichRun列表。
        style_config: 样式配置。
        level: 列表级别（0为顶级）。
    """
    text_frame = get_body_text_frame(slide)
    if text_frame is None:
        return
    
    set_autofit(text_frame)
    
    if not text_frame.paragraphs:
        text_frame.clear()
    
    for item_runs in items:
        p = text_frame.add_paragraph()
        p.level = level
        
        for run_data in item_runs:
            run = p.add_run()
            run.text = run_data.text
            
            if run_data.bold:
                run.bold = True
                apply_run_style(run, style_config.run_overrides.bold)
            
            if run_data.italic:
                run.italic = True
                apply_run_style(run, style_config.run_overrides.italic)
            
            if run_data.code:
                apply_run_style(run, style_config.run_overrides.code)


def extract_layouts(template_path: Path) -> Result[list[LayoutSpec], TemplateLoadError]:
    """从演示文稿中提取布局信息。

    参数:
        template_path: 模板文件路径。

    返回:
        Success(list[LayoutSpec])如果提取成功，Failure(TemplateLoadError)否则。
    """
    try:
        loader = TemplateLoader(template_path)
        return Success(loader.list_layouts())
    except TemplateLoadError as exc:
        return Failure(exc)
    except Exception as exc:
        return Failure(TemplateLoadError(f"提取布局失败: {exc}"))


def render_slide_item(
    slide: Slide,
    item: SlideItem,
    style_config: StyleConfig,
    renderers: dict[SlideItemType, Callable[[Slide, SlideItem, StyleConfig], None]],
) -> None:
    """根据内容类型渲染幻灯片内容项。

    参数:
        slide: Slide对象。
        item: 内容项对象。
        style_config: 样式配置。
        renderers: 渲染器注册表。
    """
    renderer = renderers.get(item.type, render_default_item)
    renderer(slide, item, style_config)


def render_default_item(slide: Slide, item: SlideItem, style_config: StyleConfig) -> None:
    """默认内容项渲染器。

    参数:
        slide: Slide对象。
        item: 内容项对象。
        style_config: 样式配置。
    """
    runs = item.meta.get("runs", [])
    
    if runs:
        render_rich_paragraph(slide, runs, style_config)
    else:
        render_paragraph(slide, item.content)


def render_image(slide: Slide, item: SlideItem, style_config: StyleConfig | None = None) -> None:
    """渲染图片内容项。

    参数:
        slide: Slide对象。
        item: 内容项对象。
        style_config: 样式配置。
    """
    prerender = item.meta.get("prerender")
    if isinstance(prerender, PrerenderResult):
        _insert_image_from_path(slide, prerender.image_path)
        return
    
    src = item.meta.get("src", "")
    if src and Path(src).exists():
        _insert_image_from_path(slide, Path(src))
    else:
        render_paragraph(slide, item.content)


def _get_image_size(image_path: Path) -> tuple[int, int]:
    """获取图片的原始尺寸（像素）。

    参数:
        image_path: 图片文件路径。

    返回:
        (宽度, 高度) 像素元组。
    """
    with PILImage.open(image_path) as img:
        return img.size


def _calculate_fit_dimensions(
    img_w_px: int,
    img_h_px: int,
    container_w: int,
    container_h: int,
) -> tuple[int, int, int, int]:
    """计算图片在容器内等比缩放后的尺寸和居中位置。

    不裁剪、不拉伸，保持原始宽高比，完整放入容器内并居中。

    参数:
        img_w_px: 图片原始宽度（像素）。
        img_h_px: 图片原始高度（像素）。
        container_w: 容器宽度（EMU）。
        container_h: 容器高度（EMU）。

    返回:
        (left, top, width, height) 均为 EMU 单位的元组。
    """
    if img_w_px == 0 or img_h_px == 0:
        return 0, 0, container_w, container_h

    img_ratio = img_w_px / img_h_px
    container_ratio = container_w / container_h

    if img_ratio > container_ratio:
        new_w = container_w
        new_h = int(container_w / img_ratio)
    else:
        new_h = container_h
        new_w = int(container_h * img_ratio)

    left = int(container_w - new_w) // 2
    top = int(container_h - new_h) // 2

    return left, top, new_w, new_h


def _find_picture_placeholder(slide: Slide) -> Any | None:
    """查找幻灯片中的PICTURE类型占位符。

    参数:
        slide: Slide对象。

    返回:
        占位符Shape对象，如果找不到则返回None。
    """
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.PICTURE:
            return shape
    return None


def _insert_image_from_path(slide: Slide, image_path: Path) -> None:
    """从文件路径插入图片到幻灯片。

    采用等比缩放+居中策略，不裁剪不拉伸。

    插入优先级:
    1. PICTURE占位符：获取其位置和尺寸，居中放置等比缩放的图片
    2. BODY文本框：获取文本框位置和尺寸，居中放置等比缩放的图片
    3. 默认：幻灯片左上角，原始尺寸

    参数:
        slide: Slide对象。
        image_path: 图片文件路径。
    """
    try:
        img_w_px, img_h_px = _get_image_size(image_path)
    except Exception:
        slide.shapes.add_picture(str(image_path), 0, 0)
        return

    picture_shape = _find_picture_placeholder(slide)
    if picture_shape is not None:
        left = picture_shape.left
        top = picture_shape.top
        width = picture_shape.width
        height = picture_shape.height
        offset_left, offset_top, new_w, new_h = _calculate_fit_dimensions(
            img_w_px, img_h_px, width, height
        )
        slide.shapes.add_picture(
            str(image_path),
            left + offset_left,
            top + offset_top,
            width=new_w,
            height=new_h,
        )
        return

    text_frame = get_body_text_frame(slide)
    if text_frame is not None:
        parent_shape = text_frame._parent
        left = parent_shape.left
        top = parent_shape.top
        width = parent_shape.width
        height = parent_shape.height
        offset_left, offset_top, new_w, new_h = _calculate_fit_dimensions(
            img_w_px, img_h_px, width, height
        )
        slide.shapes.add_picture(
            str(image_path),
            left + offset_left,
            top + offset_top,
            width=new_w,
            height=new_h,
        )
        return

    slide.shapes.add_picture(str(image_path), 0, 0)


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """将十六进制颜色转换为RGBColor。

    参数:
        hex_color: 十六进制颜色字符串，如 "#FF0000"。

    返回:
        RGBColor对象。
    """
    r, g, b = hex_to_rgb(hex_color)
    return RGBColor(r, g, b)