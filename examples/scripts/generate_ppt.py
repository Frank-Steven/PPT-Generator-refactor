"""生成样例PPT文件的脚本。

本脚本使用项目中的PPTGenerator从Markdown文件生成PPT演示文稿。
支持生成多个样例文件，并演示模板文件和主题包两种使用模式。
"""

from pathlib import Path

from pptx import Presentation

from ppt_generator import PPTGenerator
from ppt_generator.core.models import PrerenderConfig


def create_template(output_path: Path) -> None:
    """创建PPT模板文件。

    参数:
        output_path: 模板文件输出路径。
    """
    prs = Presentation()
    prs.save(str(output_path))
    print(f"  创建模板文件: {output_path.name}")


def generate_ppt_from_markdown(
    markdown_path: Path,
    template_path: Path,
    output_path: Path,
    title: str,
    use_prerender: bool = False,
) -> None:
    """从Markdown文件生成PPT。

    参数:
        markdown_path: Markdown文件路径。
        template_path: 模板文件路径。
        output_path: 输出PPT路径。
        title: 演示文稿标题。
        use_prerender: 是否启用预渲染管线。
    """
    print(f"  读取: {markdown_path.name}")

    markdown_text = markdown_path.read_text(encoding="utf-8")

    prerender_config = None
    if use_prerender:
        prerender_config = PrerenderConfig(
            enable_code=True,
            enable_mermaid=True,
            enable_latex=True,
            cache_dir=output_path.parent / ".cache" / "prerender",
            dpi=300,
            timeout=30,
        )
        print("  启用预渲染管线: 代码高亮 + Mermaid图表 + LaTeX公式")

    generator = PPTGenerator(
        markdown_text=markdown_text,
        template_path=template_path,
        output_path=output_path,
        title=title,
        prerender_config=prerender_config,
    )

    try:
        generator.generate()
        print(f"  ✓ 成功: {output_path.name}")
    except Exception as exc:
        print(f"  ✗ 失败: {exc}")
        raise


def main() -> None:
    """主函数：生成所有样例PPT。"""
    scripts_dir = Path(__file__).parent
    examples_dir = scripts_dir.parent

    template_path = examples_dir / "templates" / "template.pptx"

    if not template_path.exists():
        print("模板文件不存在，正在创建...")
        create_template(template_path)

    examples = [
        {
            "markdown": "product_intro.md",
            "output": "product_intro.pptx",
            "title": "PPT生成器 - 产品介绍",
            "use_prerender": False,
        },
        {
            "markdown": "technical_tutorial.md",
            "output": "technical_tutorial.pptx",
            "title": "Python函数式编程入门",
            "use_prerender": False,
        },
        {
            "markdown": "advanced_features.md",
            "output": "advanced_features.pptx",
            "title": "高级功能演示",
            "use_prerender": True,
        },
    ]

    print("=" * 60)
    print("开始生成样例PPT文件...")
    print("=" * 60)
    print(f"模板文件: {template_path}")
    print()

    for example in examples:
        markdown_path = examples_dir / example["markdown"]
        output_path = examples_dir / example["output"]

        if not markdown_path.exists():
            print(f"  跳过: {markdown_path.name} (文件不存在)")
            continue

        print(f"\n处理: {example['title']}")
        print(f"  {'-' * 40}")

        generate_ppt_from_markdown(
            markdown_path=markdown_path,
            template_path=template_path,
            output_path=output_path,
            title=example["title"],
            use_prerender=example["use_prerender"],
        )

    print("\n" + "=" * 60)
    print("生成完成！")
    print("=" * 60)

    verify_ppt_files(examples_dir)


def verify_ppt_files(examples_dir: Path) -> None:
    """验证生成的PPT文件。"""
    ppt_files = list(examples_dir.glob("*.pptx"))
    ppt_files = [f for f in ppt_files if f.name != "template.pptx"]

    print("\n验证结果:")
    print("-" * 60)

    for ppt_file in ppt_files:
        try:
            prs = Presentation(str(ppt_file))
            slide_count = len(prs.slides)
            print(f"  ✓ {ppt_file.name}: {slide_count} 张幻灯片")
        except Exception as exc:
            print(f"  ✗ {ppt_file.name}: 验证失败 - {exc}")


if __name__ == "__main__":
    main()
