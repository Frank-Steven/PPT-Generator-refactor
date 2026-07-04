"""创建包含所有标准布局的PPT模板。

本脚本使用python-pptx创建包含所有7种标准布局的模板文件：
1. Title Slide - 标题幻灯片
2. Title and Content - 标题和内容
3. Section Header - 章节标题
4. Two Content - 双栏内容
5. Content with Caption - 带说明的内容
6. Picture with Caption - 带说明的图片
7. Blank - 空白页
"""

from pathlib import Path

from pptx import Presentation


def create_rich_template(output_path: Path) -> None:
    """创建包含所有标准布局的PPT模板。

    参数:
        output_path: 模板文件输出路径。
    """
    prs = Presentation()

    layout_info = []
    for i, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for ph in layout.placeholders:
            ph_type = ph.placeholder_format.type
            ph_name = ph.name
            placeholders.append(f"{ph_name} ({ph_type})")

        layout_info.append({"index": i, "name": layout.name, "placeholders": placeholders})

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print("模板已创建:")
    print("=" * 60)
    for info in layout_info:
        print(f"\n布局 {info['index']}: {info['name']}")
        print(f"  占位符: {', '.join(info['placeholders'])}")
    print(f"\n文件位置: {output_path}")


def main() -> None:
    """主函数：创建模板。"""
    scripts_dir = Path(__file__).parent
    examples_dir = scripts_dir.parent
    template_path = examples_dir / "templates" / "template.pptx"

    create_rich_template(template_path)


if __name__ == "__main__":
    main()
