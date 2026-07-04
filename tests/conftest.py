"""测试配置和共享fixtures。

本文件定义了测试中使用的共享fixtures和配置，
减少测试代码重复，提高测试可维护性。
"""

from __future__ import annotations

from pathlib import Path
from typing import Generator

import pytest
from pptx import Presentation

from ppt_generator.matching import LayoutMatcher
from ppt_generator.parsers import MarkdownParser
from ppt_generator.core.models import LayoutSpec, PlaceholderSpec, SlideItem, SlideItemType, SlideSpec
from ppt_generator.templates import TemplateLoader


@pytest.fixture
def sample_markdown() -> str:
    """返回示例Markdown文本。"""
    return """# 幻灯片1

段落内容1

段落内容2

# 幻灯片2

列表项1
列表项2
"""


@pytest.fixture
def sample_markdown_with_list() -> str:
    """返回包含列表的示例Markdown文本。"""
    return """# 议程

- 介绍
- 方法论
- 结果
- 结论

# 总结

感谢观看。
"""


@pytest.fixture
def markdown_parser(sample_markdown: str) -> MarkdownParser:
    """返回MarkdownParser实例。"""
    return MarkdownParser(sample_markdown)


@pytest.fixture
def temp_pptx(tmp_path: Path) -> Path:
    """创建临时PPTX文件。"""
    pptx_path = tmp_path / "template.pptx"
    prs = Presentation()
    prs.save(str(pptx_path))
    return pptx_path


@pytest.fixture
def template_loader(temp_pptx: Path) -> TemplateLoader:
    """返回TemplateLoader实例。"""
    return TemplateLoader(temp_pptx)


@pytest.fixture
def layout_matcher() -> LayoutMatcher:
    """返回LayoutMatcher实例。"""
    return LayoutMatcher()


@pytest.fixture
def sample_slide_spec() -> SlideSpec:
    """返回示例SlideSpec实例。"""
    return SlideSpec(
        title="测试幻灯片",
        items=[
            SlideItem(type=SlideItemType.PARAGRAPH, content="测试内容"),
            SlideItem(type=SlideItemType.LIST, content="列表项"),
        ],
        layout_hint="Title Slide",
    )


@pytest.fixture
def sample_layout_spec() -> LayoutSpec:
    """返回示例LayoutSpec实例。"""
    return LayoutSpec(
        name="Title Slide",
        placeholders=[
            PlaceholderSpec(name="Title 1", placeholder_type="title", index=0, shape_id=1),
            PlaceholderSpec(name="Content Placeholder 2", placeholder_type="body", index=1, shape_id=2),
        ],
    )


@pytest.fixture
def sample_layouts(sample_layout_spec: LayoutSpec) -> list[LayoutSpec]:
    """返回示例布局列表。"""
    return [
        sample_layout_spec,
        LayoutSpec(
            name="Title and Content",
            placeholders=[
                PlaceholderSpec(name="Title 1", placeholder_type="title", index=0, shape_id=3),
                PlaceholderSpec(name="Content Placeholder 2", placeholder_type="body", index=1, shape_id=4),
            ],
        ),
        LayoutSpec(
            name="Blank",
            placeholders=[],
        ),
    ]
