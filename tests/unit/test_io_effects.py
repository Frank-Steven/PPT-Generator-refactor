"""rendering/io_effects.py 单元测试。

覆盖 IO 副作用模块的全部公共函数和关键私有辅助：
- 模板/演示文稿加载与保存（load_presentation / save_presentation / extract_layouts）
- 布局查找与幻灯片添加（find_layout_index / add_slide）
- 标题/正文渲染（render_title / get_body_text_frame / set_autofit）
- Run 样式覆盖（apply_run_style / _hex_to_rgb）
- 段落/列表/富文本渲染（render_paragraph / render_list / render_rich_paragraph / render_rich_list）
- 内容项分派（render_slide_item / render_default_item）
- 图片渲染与等比缩放（render_image / _insert_image_from_path /
  _calculate_fit_dimensions / _find_picture_placeholder / _get_image_size）
"""

from __future__ import annotations

from pathlib import Path
from typing import Any
from unittest.mock import MagicMock, patch

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
from pptx.presentation import Presentation as PresentationClass
from pptx.util import Pt
from returns.result import Failure, Result

from ppt_generator.core.exceptions import TemplateLoadError
from ppt_generator.core.models import (
    PrerenderResult,
    RichRun,
    RunStyle,
    SlideItem,
    SlideItemType,
    StyleConfig,
)
from ppt_generator.rendering import io_effects
from ppt_generator.rendering.io_effects import (
    _calculate_fit_dimensions,
    _find_first_placeholder,
    _find_picture_placeholder,
    _get_image_size,
    _hex_to_rgb,
    add_slide,
    apply_run_style,
    extract_layouts,
    find_layout_index,
    get_body_text_frame,
    load_presentation,
    render_default_item,
    render_image,
    render_list,
    render_paragraph,
    render_rich_list,
    render_rich_paragraph,
    render_slide_item,
    render_title,
    save_presentation,
    set_autofit,
)


def _is_success(result: Result[Any, Any]) -> bool:
    """判断 Result 是否为 Success（避开泛型 isinstance 限制）。"""
    return not isinstance(result, Failure)


# ──────────────────────── 共享 fixtures ────────────────────────


@pytest.fixture
def default_presentation() -> Presentation:
    """返回默认空 Presentation（自带若干内置布局）。"""
    return Presentation()


@pytest.fixture
def style_config() -> StyleConfig:
    """返回默认 StyleConfig。"""
    return StyleConfig()


@pytest.fixture
def sample_image_path(tmp_path: Path) -> Path:
    """生成一个 100x80 的 PNG 测试图片。"""
    path = tmp_path / "sample.png"
    Image.new("RGB", (100, 80), (255, 0, 0)).save(path)
    return path


def _add_title_slide(prs: Presentation) -> Any:
    """添加一个 Title Slide 布局的幻灯片。"""
    layout = next(
        (lay for lay in prs.slide_layouts if lay.name == "Title Slide"),
        prs.slide_layouts[0],
    )
    return prs.slides.add_slide(layout)


def _add_content_slide(prs: Presentation) -> Any:
    """添加一个 Title and Content 布局的幻灯片（正文占位符类型为 OBJECT）。"""
    layout = next(
        (lay for lay in prs.slide_layouts if lay.name == "Title and Content"),
        prs.slide_layouts[0],
    )
    return prs.slides.add_slide(layout)


def _add_section_header_slide(prs: Presentation) -> Any:
    """添加一个 Section Header 布局的幻灯片（正文占位符类型为 BODY）。"""
    layout = next(
        (lay for lay in prs.slide_layouts if lay.name == "Section Header"),
        prs.slide_layouts[2],
    )
    return prs.slides.add_slide(layout)


def _add_blank_slide(prs: Presentation) -> Any:
    """添加一个 Blank 布局的幻灯片。"""
    layout = next(
        (lay for lay in prs.slide_layouts if lay.name == "Blank"),
        prs.slide_layouts[-1],
    )
    return prs.slides.add_slide(layout)


# ──────────────────────── load_presentation ────────────────────────


class TestLoadPresentation:
    def test_returns_success_for_valid_template(self, tmp_path: Path) -> None:
        path = tmp_path / "template.pptx"
        Presentation().save(str(path))

        result = load_presentation(path)

        assert _is_success(result)
        assert isinstance(result.unwrap(), PresentationClass)

    def test_returns_failure_for_nonexistent_path(self, tmp_path: Path) -> None:
        path = tmp_path / "missing.pptx"

        result = load_presentation(path)

        assert isinstance(result, Failure)
        assert isinstance(result.failure(), TemplateLoadError)

    def test_returns_failure_for_corrupted_file(self, tmp_path: Path) -> None:
        path = tmp_path / "corrupt.pptx"
        path.write_text("not a pptx")

        result = load_presentation(path)

        assert isinstance(result, Failure)
        assert isinstance(result.failure(), TemplateLoadError)


# ──────────────────────── save_presentation ────────────────────────


class TestSavePresentation:
    def test_returns_success_for_valid_path(
        self, default_presentation: Presentation, tmp_path: Path
    ) -> None:
        out = tmp_path / "out.pptx"

        result = save_presentation(default_presentation, out)

        assert _is_success(result)
        assert out.exists()

    def test_creates_missing_parent_dirs(
        self, default_presentation: Presentation, tmp_path: Path
    ) -> None:
        out = tmp_path / "nested" / "deep" / "out.pptx"

        result = save_presentation(default_presentation, out)

        assert _is_success(result)
        assert out.exists()

    def test_returns_failure_when_save_fails(
        self, default_presentation: Presentation, tmp_path: Path
    ) -> None:
        # 写入到一个已存在的目录路径上（非文件），presentation.save 会抛异常
        out = tmp_path / "is_a_dir.pptx"
        out.mkdir()

        result = save_presentation(default_presentation, out)

        assert isinstance(result, Failure)
        assert isinstance(result.failure(), Exception)


# ──────────────────────── find_layout_index ────────────────────────


class TestFindLayoutIndex:
    def test_returns_success_when_layout_found(
        self, default_presentation: Presentation
    ) -> None:
        first_name = default_presentation.slide_layouts[0].name

        result = find_layout_index(default_presentation, first_name)

        assert _is_success(result)
        assert result.unwrap() == 0

    def test_returns_failure_when_layout_missing(
        self, default_presentation: Presentation
    ) -> None:
        result = find_layout_index(default_presentation, "不存在的布局")

        assert isinstance(result, Failure)
        assert isinstance(result.failure(), ValueError)
        assert "不存在的布局" in str(result.failure())


# ──────────────────────── add_slide ────────────────────────


class TestAddSlide:
    def test_adds_slide_at_index(self, default_presentation: Presentation) -> None:
        initial_count = len(default_presentation.slides._sldIdLst)

        slide = add_slide(default_presentation, 0)

        assert slide is not None
        final_count = len(default_presentation.slides._sldIdLst)
        assert final_count == initial_count + 1


# ──────────────────────── render_title ────────────────────────


class TestRenderTitle:
    def test_sets_title_text_on_title_placeholder(
        self, default_presentation: Presentation
    ) -> None:
        slide = _add_title_slide(default_presentation)

        render_title(slide, "我的标题")

        title_shape = next(
            (s for s in slide.shapes if s.is_placeholder and s.has_text_frame),
            None,
        )
        assert title_shape is not None
        assert title_shape.text == "我的标题"

    def test_no_op_when_no_title_placeholder(
        self, default_presentation: Presentation
    ) -> None:
        slide = _add_blank_slide(default_presentation)

        render_title(slide, "无标题占位符")

        # 没有占位符意味着不会抛错也不会写入文本
        titles = [
            s.text_frame.text
            for s in slide.shapes
            if s.is_placeholder and s.has_text_frame
        ]
        assert all("无标题占位符" not in t for t in titles)

    def test_falls_back_to_name_based_match(self) -> None:
        """覆盖 render_title 的名称兜底分支（L128-139）。

        构造一个 slide：第一个 next（按 TITLE/CENTER_TITLE 类型匹配）返回 None，
        第二个 next（按 "title" in name.lower() 且非 SUBTITLE 匹配）返回 mock shape。
        """
        # 构造一个非 TITLE/CENTER_TITLE/SUBTITLE 类型，但 name 含 "title" 的占位符
        title_shape = MagicMock()
        title_shape.is_placeholder = True
        title_shape.has_text_frame = True
        title_shape.name = "Custom Title Placeholder"
        title_shape.placeholder_format.type = PP_PLACEHOLDER_TYPE.OBJECT

        # 另一个干扰项：SUBTITLE 类型 + 含 "title" 名称，应被排除
        subtitle_shape = MagicMock()
        subtitle_shape.is_placeholder = True
        subtitle_shape.has_text_frame = True
        subtitle_shape.name = "Subtitle Title"
        subtitle_shape.placeholder_format.type = PP_PLACEHOLDER_TYPE.SUBTITLE

        slide = MagicMock()
        slide.shapes = [subtitle_shape, title_shape]

        render_title(slide, "fallback-title")

        # 应该写入到 title_shape（OBJECT 类型 + 含 title 名称）
        assert title_shape.text == "fallback-title"
        # subtitle_shape 不应被写入
        assert not hasattr(subtitle_shape, "text") or subtitle_shape.text != "fallback-title"


# ──────────────────────── _find_first_placeholder / get_body_text_frame ────────────────────────


class TestFindFirstPlaceholder:
    def test_returns_text_frame_for_body(
        self, default_presentation: Presentation
    ) -> None:
        # Section Header 布局有 BODY 占位符
        slide = _add_section_header_slide(default_presentation)

        tf = _find_first_placeholder(
            slide, (PP_PLACEHOLDER_TYPE.BODY,)
        )

        assert tf is not None
        assert hasattr(tf, "add_paragraph")

    def test_returns_text_frame_for_object(
        self, default_presentation: Presentation
    ) -> None:
        # Title and Content 布局用 OBJECT 占位符
        slide = _add_content_slide(default_presentation)

        tf = _find_first_placeholder(
            slide, (PP_PLACEHOLDER_TYPE.OBJECT,)
        )

        assert tf is not None

    def test_returns_none_when_no_matching_type(
        self, default_presentation: Presentation
    ) -> None:
        slide = _add_blank_slide(default_presentation)

        tf = _find_first_placeholder(slide, (PP_PLACEHOLDER_TYPE.BODY,))

        assert tf is None

    def test_falls_through_to_second_priority_type(
        self, default_presentation: Presentation
    ) -> None:
        # Title Slide 的副标题是 SUBTITLE，BODY/OBJECT 不存在
        slide = _add_title_slide(default_presentation)

        tf = _find_first_placeholder(
            slide,
            (PP_PLACEHOLDER_TYPE.BODY, PP_PLACEHOLDER_TYPE.SUBTITLE),
        )

        assert tf is not None


class TestGetBodyTextFrame:
    def test_returns_object_for_content_layout(
        self, default_presentation: Presentation
    ) -> None:
        # Title and Content 用 OBJECT（get_body_text_frame 第二优先级）
        slide = _add_content_slide(default_presentation)

        tf = get_body_text_frame(slide)

        assert tf is not None

    def test_returns_body_for_section_header(
        self, default_presentation: Presentation
    ) -> None:
        # Section Header 用 BODY（最高优先级）
        slide = _add_section_header_slide(default_presentation)

        tf = get_body_text_frame(slide)

        assert tf is not None

    def test_returns_subtitle_for_title_slide(
        self, default_presentation: Presentation
    ) -> None:
        # Title Slide 只有 SUBTITLE（最低优先级兜底）
        slide = _add_title_slide(default_presentation)

        tf = get_body_text_frame(slide)

        assert tf is not None

    def test_returns_none_for_blank_layout(
        self, default_presentation: Presentation
    ) -> None:
        slide = _add_blank_slide(default_presentation)

        tf = get_body_text_frame(slide)

        assert tf is None


# ──────────────────────── set_autofit ────────────────────────


class TestSetAutofit:
    def test_sets_auto_size_on_text_frame(
        self, default_presentation: Presentation
    ) -> None:
        slide = _add_content_slide(default_presentation)
        tf = get_body_text_frame(slide)
        assert tf is not None

        set_autofit(tf)

        # MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE 应被设置
        from pptx.enum.text import MSO_AUTO_SIZE

        assert tf.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


# ──────────────────────── apply_run_style / _hex_to_rgb ────────────────────────


class _SpyFont:
    """记录所有属性赋值的 Font 替身。"""

    # 声明所有可能被 apply_run_style 设置的属性，避免 mypy attr-defined
    name: Any
    size: Any
    underline: Any
    highlight_color: Any

    def __init__(self) -> None:
        self.set_attrs: set[str] = set()
        self.color = _SpyColor()

    def __setattr__(self, name: str, value: Any) -> None:
        if name not in ("set_attrs", "color"):
            self.set_attrs.add(name)
        super().__setattr__(name, value)


class _SpyColor:
    rgb: Any

    def __init__(self) -> None:
        self.rgb: Any = None

    def __setattr__(self, name: str, value: Any) -> None:
        super().__setattr__(name, value)


class _SpyRun:
    """记录所有属性赋值的 Run 替身。"""

    # 声明所有可能被 apply_run_style 设置的属性
    bold: Any
    italic: Any

    def __init__(self) -> None:
        self.set_attrs: set[str] = set()
        self.font = _SpyFont()

    def __setattr__(self, name: str, value: Any) -> None:
        if name not in ("set_attrs", "font"):
            self.set_attrs.add(name)
        super().__setattr__(name, value)


class TestApplyRunStyle:
    def test_applies_all_fields(self) -> None:
        run = _SpyRun()

        style = RunStyle(
            font="Arial",
            font_size=20,
            color="#FF0000",
            bold=True,
            italic=True,
            underline=True,
            background_color="#00FF00",
        )

        apply_run_style(run, style)

        assert run.font.name == "Arial"
        assert run.font.size == Pt(20)
        assert run.bold is True
        assert run.italic is True
        assert run.font.underline is True
        assert run.font.color.rgb is not None
        assert run.font.highlight_color is not None
        # 应设置全部 7 个字段
        assert {"name", "size", "underline", "highlight_color"} <= run.font.set_attrs

    def test_applies_nothing_when_all_none(self) -> None:
        run = _SpyRun()

        apply_run_style(run, RunStyle())

        # run 上不应有任何属性被赋值（bold/italic 都未被设置）
        assert run.set_attrs == set()
        assert run.font.set_attrs == set()

    def test_applies_only_font_when_only_font_set(self) -> None:
        run = _SpyRun()

        apply_run_style(run, RunStyle(font="Courier"))

        assert run.font.name == "Courier"
        assert run.font.set_attrs == {"name"}
        assert run.set_attrs == set()


class TestHexToRgb:
    def test_parses_hex_with_hash(self) -> None:
        rgb = _hex_to_rgb("#FF0000")
        assert (rgb[0], rgb[1], rgb[2]) == (255, 0, 0)

    def test_parses_hex_without_hash(self) -> None:
        rgb = _hex_to_rgb("00FF00")
        assert (rgb[0], rgb[1], rgb[2]) == (0, 255, 0)


# ──────────────────────── render_paragraph ────────────────────────


class TestRenderParagraph:
    def test_appends_paragraph_when_append_true(
        self, default_presentation: Presentation
    ) -> None:
        slide = _add_content_slide(default_presentation)
        render_paragraph(slide, "first", append=True)
        render_paragraph(slide, "second", append=True)

        tf = get_body_text_frame(slide)
        texts = [p.text for p in tf.paragraphs if p.text]
        assert "first" in texts
        assert "second" in texts

    def test_clears_when_append_false(
        self, default_presentation: Presentation
    ) -> None:
        slide = _add_content_slide(default_presentation)
        render_paragraph(slide, "first", append=True)
        render_paragraph(slide, "second", append=False)

        tf = get_body_text_frame(slide)
        texts = [p.text for p in tf.paragraphs if p.text]
        assert "first" not in texts
        assert "second" in texts

    def test_no_op_when_no_body_text_frame(
        self, default_presentation: Presentation
    ) -> None:
        slide = _add_blank_slide(default_presentation)
        # 不应抛错
        render_paragraph(slide, "anything")


# ──────────────────────── render_list ────────────────────────


class TestRenderList:
    def test_adds_list_item_with_level(
        self, default_presentation: Presentation
    ) -> None:
        slide = _add_content_slide(default_presentation)

        render_list(slide, "item1", level=0)
        render_list(slide, "item2", level=1)

        tf = get_body_text_frame(slide)
        # 找到包含 item1 和 item2 的段落
        item1_p = next(p for p in tf.paragraphs if p.text == "item1")
        item2_p = next(p for p in tf.paragraphs if p.text == "item2")
        assert item1_p.level == 0
        assert item2_p.level == 1

    def test_clears_when_text_frame_has_no_paragraphs(self) -> None:
        """覆盖 L303-304 分支：text_frame.paragraphs 为空时调用 clear()。"""
        tf = MagicMock()
        tf.paragraphs = []  # 触发 clear() 分支
        added_p = MagicMock()
        tf.add_paragraph.return_value = added_p

        slide = MagicMock()
        # 让 get_body_text_frame 返回我们的 mock
        with patch.object(io_effects, "get_body_text_frame", return_value=tf):
            render_list(slide, "lonely-item", level=0)

        tf.clear.assert_called_once()
        tf.add_paragraph.assert_called_once()
        assert added_p.text == "lonely-item"
        assert added_p.level == 0

    def test_no_op_when_no_body_text_frame(
        self, default_presentation: Presentation
    ) -> None:
        slide = _add_blank_slide(default_presentation)
        render_list(slide, "item")


# ──────────────────────── render_rich_paragraph ────────────────────────


class TestRenderRichParagraph:
    def test_renders_runs_with_styles(
        self, default_presentation: Presentation, style_config: StyleConfig
    ) -> None:
        slide = _add_content_slide(default_presentation)
        runs = [
            RichRun(text="普通", bold=False, italic=False, code=False),
            RichRun(text="粗体", bold=True),
            RichRun(text="斜体", italic=True),
            RichRun(text="代码", code=True),
            RichRun(text="链接", link="https://example.com"),
            RichRun(text="删除", strikethrough=True),
        ]

        render_rich_paragraph(slide, runs, style_config, append=True)

        tf = get_body_text_frame(slide)
        all_text = "".join(r.text for p in tf.paragraphs for r in p.runs)
        assert "普通" in all_text
        assert "粗体" in all_text
        assert "斜体" in all_text
        assert "代码" in all_text
        assert "链接" in all_text
        assert "删除" in all_text

    def test_clears_when_append_false(
        self, default_presentation: Presentation, style_config: StyleConfig
    ) -> None:
        slide = _add_content_slide(default_presentation)
        render_rich_paragraph(
            slide, [RichRun(text="first")], style_config, append=True
        )
        render_rich_paragraph(
            slide, [RichRun(text="second")], style_config, append=False
        )

        tf = get_body_text_frame(slide)
        all_text = "".join(r.text for p in tf.paragraphs for r in p.runs)
        assert "first" not in all_text
        assert "second" in all_text

    def test_no_op_when_no_body_text_frame(
        self, default_presentation: Presentation, style_config: StyleConfig
    ) -> None:
        slide = _add_blank_slide(default_presentation)
        render_rich_paragraph(slide, [RichRun(text="x")], style_config)


# ──────────────────────── render_rich_list ────────────────────────


class TestRenderRichList:
    def test_renders_multiple_items(
        self, default_presentation: Presentation, style_config: StyleConfig
    ) -> None:
        slide = _add_content_slide(default_presentation)
        items = [
            [RichRun(text="item1-bold", bold=True)],
            [RichRun(text="item2-italic", italic=True)],
            [RichRun(text="item3-code", code=True)],
        ]

        render_rich_list(slide, items, style_config, level=1)

        tf = get_body_text_frame(slide)
        all_text = "".join(r.text for p in tf.paragraphs for r in p.runs)
        assert "item1-bold" in all_text
        assert "item2-italic" in all_text
        assert "item3-code" in all_text

    def test_clears_when_text_frame_has_no_paragraphs(
        self, style_config: StyleConfig
    ) -> None:
        """覆盖 L331-332 分支：text_frame.paragraphs 为空时调用 clear()。"""
        tf = MagicMock()
        tf.paragraphs = []  # 触发 clear() 分支
        added_p = MagicMock()
        added_p.runs = []
        tf.add_paragraph.return_value = added_p

        slide = MagicMock()
        items = [[RichRun(text="x")]]

        with patch.object(io_effects, "get_body_text_frame", return_value=tf):
            render_rich_list(slide, items, style_config, level=0)

        tf.clear.assert_called_once()
        # items 有 1 项，应调用 1 次 add_paragraph
        assert tf.add_paragraph.call_count == 1

    def test_no_op_when_no_body_text_frame(
        self, default_presentation: Presentation, style_config: StyleConfig
    ) -> None:
        slide = _add_blank_slide(default_presentation)
        render_rich_list(slide, [[RichRun(text="x")]], style_config)


# ──────────────────────── extract_layouts ────────────────────────


class TestExtractLayouts:
    def test_returns_success_for_valid_template(self, tmp_path: Path) -> None:
        path = tmp_path / "template.pptx"
        Presentation().save(str(path))

        result = extract_layouts(path)

        assert _is_success(result)
        layouts = result.unwrap()
        assert isinstance(layouts, list)
        assert len(layouts) > 0

    def test_returns_failure_for_missing_file(self, tmp_path: Path) -> None:
        path = tmp_path / "missing.pptx"

        result = extract_layouts(path)

        assert isinstance(result, Failure)
        assert isinstance(result.failure(), TemplateLoadError)

    def test_returns_failure_for_corrupted_file(self, tmp_path: Path) -> None:
        path = tmp_path / "corrupt.pptx"
        path.write_text("not a pptx")

        result = extract_layouts(path)

        assert isinstance(result, Failure)
        assert isinstance(result.failure(), TemplateLoadError)


# ──────────────────────── render_slide_item / render_default_item ────────────────────────


class TestRenderSlideItem:
    def test_dispatches_to_registered_renderer(
        self, default_presentation: Presentation, style_config: StyleConfig
    ) -> None:
        slide = _add_content_slide(default_presentation)
        called: list[bool] = []

        def custom_renderer(
            s: Any, item: SlideItem, cfg: StyleConfig
        ) -> None:
            called.append(True)

        renderers = {SlideItemType.PARAGRAPH: custom_renderer}
        item = SlideItem(type=SlideItemType.PARAGRAPH, content="x")

        render_slide_item(slide, item, style_config, renderers)

        assert called == [True]

    def test_falls_back_to_default_when_type_not_registered(
        self, default_presentation: Presentation, style_config: StyleConfig
    ) -> None:
        slide = _add_content_slide(default_presentation)
        renderers: dict[SlideItemType, Any] = {}
        item = SlideItem(type=SlideItemType.PARAGRAPH, content="fallback-text")

        render_slide_item(slide, item, style_config, renderers)

        tf = get_body_text_frame(slide)
        all_text = "".join(p.text for p in tf.paragraphs)
        assert "fallback-text" in all_text


class TestRenderDefaultItem:
    def test_uses_rich_runs_when_present(
        self, default_presentation: Presentation, style_config: StyleConfig
    ) -> None:
        slide = _add_content_slide(default_presentation)
        item = SlideItem(
            type=SlideItemType.PARAGRAPH,
            content="plain",
            meta={"runs": [RichRun(text="rich-text", bold=True)]},
        )

        render_default_item(slide, item, style_config)

        tf = get_body_text_frame(slide)
        all_text = "".join(r.text for p in tf.paragraphs for r in p.runs)
        assert "rich-text" in all_text

    def test_falls_back_to_plain_content(
        self, default_presentation: Presentation, style_config: StyleConfig
    ) -> None:
        slide = _add_content_slide(default_presentation)
        item = SlideItem(type=SlideItemType.PARAGRAPH, content="plain-only")

        render_default_item(slide, item, style_config)

        tf = get_body_text_frame(slide)
        all_text = "".join(p.text for p in tf.paragraphs)
        assert "plain-only" in all_text


# ──────────────────────── render_image ────────────────────────


class TestRenderImage:
    def test_inserts_image_from_prerender(
        self,
        default_presentation: Presentation,
        style_config: StyleConfig,
        sample_image_path: Path,
    ) -> None:
        slide = _add_content_slide(default_presentation)
        prerender = PrerenderResult(
            image_path=sample_image_path,
            width=100,
            height=80,
            content_hash="abc",
        )
        item = SlideItem(
            type=SlideItemType.IMAGE,
            content="alt",
            meta={"prerender": prerender},
        )

        render_image(slide, item, style_config)

        pictures = [s for s in slide.shapes if s.shape_type == 13]  # PICTURE
        assert len(pictures) >= 1

    def test_inserts_image_from_src_when_exists(
        self,
        default_presentation: Presentation,
        style_config: StyleConfig,
        sample_image_path: Path,
    ) -> None:
        slide = _add_content_slide(default_presentation)
        item = SlideItem(
            type=SlideItemType.IMAGE,
            content="alt",
            meta={"src": str(sample_image_path)},
        )

        render_image(slide, item, style_config)

        pictures = [s for s in slide.shapes if s.shape_type == 13]
        assert len(pictures) >= 1

    def test_renders_paragraph_when_src_missing(
        self,
        default_presentation: Presentation,
        style_config: StyleConfig,
    ) -> None:
        slide = _add_content_slide(default_presentation)
        item = SlideItem(
            type=SlideItemType.IMAGE,
            content="image-not-found-text",
            meta={"src": "/nonexistent/path.png"},
        )

        render_image(slide, item, style_config)

        tf = get_body_text_frame(slide)
        all_text = "".join(p.text for p in tf.paragraphs)
        assert "image-not-found-text" in all_text

    def test_renders_paragraph_when_no_src(
        self,
        default_presentation: Presentation,
        style_config: StyleConfig,
    ) -> None:
        slide = _add_content_slide(default_presentation)
        item = SlideItem(type=SlideItemType.IMAGE, content="no-src-text")

        render_image(slide, item, style_config)

        tf = get_body_text_frame(slide)
        all_text = "".join(p.text for p in tf.paragraphs)
        assert "no-src-text" in all_text

    def test_style_config_optional(
        self,
        default_presentation: Presentation,
        sample_image_path: Path,
    ) -> None:
        slide = _add_content_slide(default_presentation)
        item = SlideItem(
            type=SlideItemType.IMAGE,
            content="alt",
            meta={"src": str(sample_image_path)},
        )

        # 不传 style_config 也应能工作
        render_image(slide, item)

        pictures = [s for s in slide.shapes if s.shape_type == 13]
        assert len(pictures) >= 1


# ──────────────────────── _get_image_size ────────────────────────


class TestGetImageSize:
    def test_returns_width_height(self, sample_image_path: Path) -> None:
        w, h = _get_image_size(sample_image_path)

        assert w == 100
        assert h == 80


# ──────────────────────── _calculate_fit_dimensions ────────────────────────


class TestCalculateFitDimensions:
    def test_uses_full_container_when_image_wider(self) -> None:
        # img_ratio = 2.0, container_ratio = 1.0 → 宽图，按宽缩放
        left, top, w, h = _calculate_fit_dimensions(
            img_w_px=200, img_h_px=100, container_w=1000, container_h=1000
        )
        assert w == 1000
        assert h == 500
        assert left == 0
        assert top == 250

    def test_uses_full_container_when_image_taller(self) -> None:
        # img_ratio = 0.5, container_ratio = 1.0 → 高图，按高缩放
        left, top, w, h = _calculate_fit_dimensions(
            img_w_px=100, img_h_px=200, container_w=1000, container_h=1000
        )
        assert h == 1000
        assert w == 500
        assert left == 250
        assert top == 0

    def test_returns_full_container_when_image_zero(self) -> None:
        left, top, w, h = _calculate_fit_dimensions(
            img_w_px=0, img_h_px=0, container_w=500, container_h=400
        )
        assert (left, top, w, h) == (0, 0, 500, 400)

    def test_perfect_fit_no_offset(self) -> None:
        # img_ratio == container_ratio，取宽图分支但结果一致
        left, top, w, h = _calculate_fit_dimensions(
            img_w_px=200, img_h_px=200, container_w=800, container_h=800
        )
        assert w == 800
        assert h == 800
        assert left == 0
        assert top == 0


# ──────────────────────── _find_picture_placeholder ────────────────────────


class TestFindPicturePlaceholder:
    def test_returns_none_when_no_picture_placeholder(
        self, default_presentation: Presentation
    ) -> None:
        slide = _add_content_slide(default_presentation)

        result = _find_picture_placeholder(slide)

        assert result is None

    def test_returns_shape_when_picture_placeholder_exists(
        self, default_presentation: Presentation
    ) -> None:
        # 找到含 PICTURE 占位符的布局（默认模板的 "Picture with Caption"）
        pic_layout = next(
            (
                lay
                for lay in default_presentation.slide_layouts
                if any(
                    s.is_placeholder
                    and s.placeholder_format.type == PP_PLACEHOLDER_TYPE.PICTURE
                    for s in lay.placeholders
                )
            ),
            None,
        )
        if pic_layout is None:
            pytest.skip("默认模板不含 PICTURE 占位符布局")
        slide = default_presentation.slides.add_slide(pic_layout)

        result = _find_picture_placeholder(slide)

        assert result is not None


# ──────────────────────── _insert_image_from_path ────────────────────────


class TestInsertImageFromPath:
    def test_inserts_into_picture_placeholder_when_present(
        self,
        default_presentation: Presentation,
        sample_image_path: Path,
    ) -> None:
        pic_layout = next(
            (
                lay
                for lay in default_presentation.slide_layouts
                if any(
                    s.is_placeholder
                    and s.placeholder_format.type == PP_PLACEHOLDER_TYPE.PICTURE
                    for s in lay.placeholders
                )
            ),
            None,
        )
        if pic_layout is None:
            pytest.skip("默认模板不含 PICTURE 占位符布局")
        slide = default_presentation.slides.add_slide(pic_layout)

        from ppt_generator.rendering.io_effects import _insert_image_from_path

        _insert_image_from_path(slide, sample_image_path)

        pictures = [s for s in slide.shapes if s.shape_type == 13]
        assert len(pictures) >= 1

    def test_inserts_into_body_text_frame_when_no_picture_placeholder(
        self,
        default_presentation: Presentation,
        sample_image_path: Path,
    ) -> None:
        slide = _add_content_slide(default_presentation)

        from ppt_generator.rendering.io_effects import _insert_image_from_path

        _insert_image_from_path(slide, sample_image_path)

        pictures = [s for s in slide.shapes if s.shape_type == 13]
        assert len(pictures) >= 1

    def test_inserts_at_origin_when_no_text_frame(
        self,
        default_presentation: Presentation,
        sample_image_path: Path,
    ) -> None:
        slide = _add_blank_slide(default_presentation)

        from ppt_generator.rendering.io_effects import _insert_image_from_path

        _insert_image_from_path(slide, sample_image_path)

        pictures = [s for s in slide.shapes if s.shape_type == 13]
        assert len(pictures) >= 1
        # 没有占位符，落在 (0, 0)
        assert pictures[0].left == 0
        assert pictures[0].top == 0

    def test_falls_back_to_origin_when_size_unreadable(
        self,
        default_presentation: Presentation,
        tmp_path: Path,
    ) -> None:
        slide = _add_content_slide(default_presentation)
        bad_path = tmp_path / "not-an-image.png"
        bad_path.write_text("garbage")

        from ppt_generator.rendering.io_effects import _insert_image_from_path

        # _get_image_size 抛异常 → 走 except 分支直接 add_picture(0, 0)
        # mock slide.shapes.add_picture 避免底层 PIL 再次抛错
        slide.shapes.add_picture = MagicMock()

        _insert_image_from_path(slide, bad_path)

        # 验证走了 except 分支，调用 add_picture(str(bad_path), 0, 0)
        slide.shapes.add_picture.assert_called_once_with(str(bad_path), 0, 0)


# ──────────────────────── 模块结构 ────────────────────────


class TestModuleStructure:
    def test_render_default_item_is_default_exported(self) -> None:
        """render_default_item 应是模块级别可访问的（render_slide_item 用作默认）。"""
        assert hasattr(io_effects, "render_default_item")
        assert callable(io_effects.render_default_item)
