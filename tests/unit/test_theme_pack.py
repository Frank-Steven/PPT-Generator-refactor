"""主题包加载器单元测试。

覆盖 load_theme_pack 主入口、manifest/style/layouts 解析、
模板验证、错误路径，以及 list_available_themes。
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from ppt_generator.core.exceptions import InvalidConfigError, MissingFileError, TemplateLoadError
from ppt_generator.core.models import (
    LayoutConfig,
    LayoutDefaults,
    StyleConfig,
    ThemePackManifest,
)
from ppt_generator.themes.theme_pack import (
    _load_layout_config,
    _load_manifest,
    _load_style_config,
    _parse_layout_config,
    _parse_layout_def,
    _parse_run_overrides,
    _parse_run_style,
    _parse_style_config,
    _resolve_optional_path,
    _resolve_template_path,
    _validate_template,
    list_available_themes,
    load_theme_pack,
)

# ──────────────────────── 测试用主题包构造 helper ────────────────────────


def _create_pptx_with_layouts(path: Path, layout_names: list[str]) -> None:
    """创建一个含指定布局名称的 pptx 文件。

    python-pptx 默认创建的 Presentation 含 11 个内建布局，名称固定。
    这里直接使用默认布局名。
    """
    prs = Presentation()
    # 默认布局名: Title Slide, Title and Content, Section Header, Two Content,
    # Comparison, Title Only, Blank, Content with Caption, Picture with Caption, Title and Vertical Text, Vertical Title and Text
    # 通过修改 layout.name 来匹配测试需求
    for i, name in enumerate(layout_names):
        if i < len(prs.slide_layouts):
            prs.slide_layouts[i].name = name
    prs.save(str(path))


@pytest.fixture
def valid_theme_pack(tmp_path: Path) -> Path:
    """创建一个完整合法的主题包目录。"""
    pack_dir = tmp_path / "test_theme"
    pack_dir.mkdir()

    # manifest.yaml
    (pack_dir / "manifest.yaml").write_text(
        """
name: 测试主题
version: 1.0.0
author: 测试者
description: 测试用主题包
spec_version: "1.0"
files:
  template: template.pptx
  style: style.yaml
  layouts: layouts.yaml
  preview: preview.png
preview:
  path: preview.png
tags:
  - 测试
  - 示例
""",
        encoding="utf-8",
    )

    # style.yaml
    (pack_dir / "style.yaml").write_text(
        """
code:
  font: Consolas
  font_size: 12
  theme: monokai
  line_numbers: true
  background_color: "#272822"
  text_color: "#F8F8F2"
  border_radius: 4
  padding: 12
  line_height: 1.4
mermaid:
  theme: dark
  background_color: "#1a1a1a"
  scale: 2
  padding: 10
latex:
  font_size: 14
  background_color: transparent
  dpi: 300
  color: "#333333"
table:
  font: 微软雅黑
  font_size: 10
  header_background: "#4472C4"
  header_color: "#FFFFFF"
  even_row_background: "#F5F5F5"
  odd_row_background: "#FFFFFF"
  border_color: "#CCCCCC"
  border_width: 1
run_overrides:
  bold:
    font: Arial
    color: "#FF0000"
    bold: true
""",
        encoding="utf-8",
    )

    # layouts.yaml - 名称与 template.pptx 中的布局一致
    (pack_dir / "layouts.yaml").write_text(
        """
version: "1.0"
defaults:
  default: title-and-content
  first_slide: title-slide
  section_divider: section-header
  content: title-and-content
  multi_column: two-content
  media: content-with-caption
  image: picture-with-caption
  full_width: blank
groups:
  content:
    display_name: 内容布局
    description: 内容相关布局
layouts:
  - id: title-slide
    name: Title Slide
    display_name: 标题幻灯片
    description: 标题页
    group: content
    placeholders:
      - index: 0
        type: title
        role: title
      - index: 1
        type: body
        role: subtitle
    keywords:
      - 标题
      - 首页
    tags:
      - basic
    auto_apply:
      conditions:
        - first_slide
      priority: 100
  - id: title-and-content
    name: Title and Content
    display_name: 标题和内容
    placeholders:
      - index: 0
        type: title
      - index: 1
        type: body
""",
        encoding="utf-8",
    )

    # template.pptx - 布局名匹配 layouts.yaml 中定义
    _create_pptx_with_layouts(pack_dir / "template.pptx", ["Title Slide", "Title and Content"])

    # preview.png 可选文件（不存在也行，这里创建一个空文件以测试 _resolve_optional_path）
    (pack_dir / "preview.png").write_bytes(b"\x89PNG\r\n\x1a\n")

    return pack_dir


# ──────────────────────── load_theme_pack 主入口测试 ────────────────────────


class TestLoadThemePack:
    """测试 load_theme_pack 主入口。"""

    def test_loads_valid_theme_pack(self, valid_theme_pack: Path) -> None:
        result = load_theme_pack(valid_theme_pack)

        assert result.manifest.name == "测试主题"
        assert result.manifest.version == "1.0.0"
        assert result.manifest.author == "测试者"
        assert result.template_path.exists()
        assert isinstance(result.style_config, StyleConfig)
        assert result.layout_config.version == "1.0"
        assert len(result.layout_config.layouts) == 2
        assert result.preview_path is not None
        assert result.preview_path.exists()

    def test_nonexistent_path_raises_missing_file_error(self, tmp_path: Path) -> None:
        with pytest.raises(MissingFileError):
            load_theme_pack(tmp_path / "nonexistent")

    def test_file_path_raises_invalid_config_error(self, tmp_path: Path) -> None:
        """传入文件而非目录应抛 InvalidConfigError。"""
        file_path = tmp_path / "not_a_dir.txt"
        file_path.touch()
        with pytest.raises(InvalidConfigError):
            load_theme_pack(file_path)

    def test_missing_manifest_raises_missing_file_error(self, tmp_path: Path) -> None:
        pack = tmp_path / "no_manifest"
        pack.mkdir()
        with pytest.raises(FileNotFoundError):
            load_theme_pack(pack)

    def test_missing_template_raises_missing_file_error(self, tmp_path: Path) -> None:
        """manifest 中指定的模板不存在应抛 MissingFileError。"""
        pack = tmp_path / "no_template"
        pack.mkdir()
        (pack / "manifest.yaml").write_text(
            "name: t\nversion: 1.0\nauthor: a\nfiles:\n  template: missing.pptx\n",
            encoding="utf-8",
        )
        (pack / "layouts.yaml").write_text(
            "version: '1.0'\nlayouts: []\n", encoding="utf-8"
        )
        with pytest.raises(MissingFileError):
            load_theme_pack(pack)

    def test_missing_layouts_raises_missing_file_error(self, tmp_path: Path) -> None:
        pack = tmp_path / "no_layouts"
        pack.mkdir()
        (pack / "manifest.yaml").write_text(
            "name: t\nversion: 1.0\nauthor: a\n", encoding="utf-8"
        )
        (pack / "template.pptx").touch()  # 空文件，但 _validate_template 会先检查 layouts
        with pytest.raises(MissingFileError):
            load_theme_pack(pack)

    def test_accepts_string_path(self, valid_theme_pack: Path) -> None:
        """load_theme_pack 应接受 str 路径。"""
        result = load_theme_pack(str(valid_theme_pack))
        assert result.manifest.name == "测试主题"

    def test_fonts_and_assets_dirs_detected(self, valid_theme_pack: Path) -> None:
        """fonts/ 和 assets/ 目录存在时被识别。"""
        (valid_theme_pack / "fonts").mkdir()
        (valid_theme_pack / "fonts" / "custom.ttf").touch()
        (valid_theme_pack / "assets").mkdir()
        (valid_theme_pack / "assets" / "logo.png").touch()

        result = load_theme_pack(valid_theme_pack)
        assert result.fonts_path is not None
        assert result.fonts_path.exists()
        assert result.assets_path is not None
        assert result.assets_path.exists()

    def test_fonts_and_assets_dirs_none_when_absent(self, valid_theme_pack: Path) -> None:
        """fonts/ 和 assets/ 目录不存在时为 None。"""
        result = load_theme_pack(valid_theme_pack)
        assert result.fonts_path is None
        assert result.assets_path is None


# ──────────────────────── manifest 解析测试 ────────────────────────


class TestLoadManifest:
    """测试 _load_manifest。"""

    def test_loads_complete_manifest(self, valid_theme_pack: Path) -> None:
        manifest = _load_manifest(valid_theme_pack)
        assert manifest.name == "测试主题"
        assert manifest.version == "1.0.0"
        assert manifest.author == "测试者"
        assert manifest.description == "测试用主题包"
        assert manifest.spec_version == "1.0"
        assert manifest.files["template"] == "template.pptx"
        assert "测试" in manifest.tags

    def test_uses_default_spec_version_when_absent(self, tmp_path: Path) -> None:
        (tmp_path / "manifest.yaml").write_text(
            "name: t\nversion: 1.0\nauthor: a\n", encoding="utf-8"
        )
        manifest = _load_manifest(tmp_path)
        assert manifest.spec_version == "1.0"

    def test_empty_name_raises_invalid_config(self, tmp_path: Path) -> None:
        (tmp_path / "manifest.yaml").write_text(
            "name: ''\nversion: 1.0\nauthor: a\n", encoding="utf-8"
        )
        with pytest.raises(InvalidConfigError):
            _load_manifest(tmp_path)

    def test_empty_version_raises_invalid_config(self, tmp_path: Path) -> None:
        (tmp_path / "manifest.yaml").write_text(
            "name: t\nversion: ''\nauthor: a\n", encoding="utf-8"
        )
        with pytest.raises(InvalidConfigError):
            _load_manifest(tmp_path)

    def test_empty_author_raises_invalid_config(self, tmp_path: Path) -> None:
        (tmp_path / "manifest.yaml").write_text(
            "name: t\nversion: 1.0\nauthor: ''\n", encoding="utf-8"
        )
        with pytest.raises(InvalidConfigError):
            _load_manifest(tmp_path)


# ──────────────────────── 路径解析测试 ────────────────────────


class TestResolveTemplatePath:
    """测试 _resolve_template_path。"""

    def test_resolves_existing_template(self, valid_theme_pack: Path) -> None:
        manifest = _load_manifest(valid_theme_pack)
        path = _resolve_template_path(valid_theme_pack, manifest)
        assert path == valid_theme_pack / "template.pptx"
        assert path.exists()

    def test_uses_default_filename_when_not_specified(self, tmp_path: Path) -> None:
        """manifest 未指定 template 时默认使用 template.pptx。"""
        (tmp_path / "template.pptx").touch()
        manifest = ThemePackManifest(name="t", version="1.0", author="a", files={})
        path = _resolve_template_path(tmp_path, manifest)
        assert path == tmp_path / "template.pptx"

    def test_missing_template_raises_missing_file(self, tmp_path: Path) -> None:
        manifest = ThemePackManifest(
            name="t", version="1.0", author="a", files={"template": "missing.pptx"}
        )
        with pytest.raises(MissingFileError):
            _resolve_template_path(tmp_path, manifest)


class TestResolveOptionalPath:
    """测试 _resolve_optional_path。"""

    def test_none_filename_returns_none(self, tmp_path: Path) -> None:
        assert _resolve_optional_path(tmp_path, None) is None

    def test_empty_filename_returns_none(self, tmp_path: Path) -> None:
        assert _resolve_optional_path(tmp_path, "") is None

    def test_existing_file_returns_path(self, tmp_path: Path) -> None:
        (tmp_path / "preview.png").touch()
        path = _resolve_optional_path(tmp_path, "preview.png")
        assert path is not None
        assert path.exists()

    def test_nonexistent_file_returns_none(self, tmp_path: Path) -> None:
        path = _resolve_optional_path(tmp_path, "missing.png")
        assert path is None


# ──────────────────────── style 配置解析测试 ────────────────────────


class TestLoadStyleConfig:
    """测试 _load_style_config。"""

    def test_loads_complete_style(self, valid_theme_pack: Path) -> None:
        manifest = _load_manifest(valid_theme_pack)
        style = _load_style_config(valid_theme_pack, manifest)

        assert style.code.font == "Consolas"
        assert style.code.font_size == 12
        assert style.code.theme == "monokai"
        assert style.mermaid.theme == "dark"
        assert style.latex.font_size == 14
        assert style.table.font == "微软雅黑"
        assert style.run_overrides.bold.color == "#FF0000"

    def test_missing_style_file_returns_default(self, tmp_path: Path) -> None:
        """style 文件不存在时返回默认 StyleConfig。"""
        manifest = ThemePackManifest(
            name="t", version="1.0", author="a",
            files={"style": "nonexistent.yaml"},
        )
        style = _load_style_config(tmp_path, manifest)
        assert isinstance(style, StyleConfig)
        # 验证使用了默认值
        assert style.code.font == "Consolas"


class TestParseStyleConfig:
    """测试 _parse_style_config 直接解析字典。"""

    def test_empty_dict_uses_all_defaults(self) -> None:
        style = _parse_style_config({})
        assert isinstance(style, StyleConfig)
        assert style.code.font == "Consolas"
        assert style.mermaid.theme == "dark"
        assert style.latex.font_size == 14
        assert style.table.font == "微软雅黑"

    def test_partial_dict_preserves_defaults(self) -> None:
        style = _parse_style_config({"code": {"font_size": 20}})
        assert style.code.font_size == 20
        # 未指定的字段使用默认值
        assert style.code.font == "Consolas"

    def test_run_overrides_parsed(self) -> None:
        style = _parse_style_config({
            "run_overrides": {
                "italic": {"color": "#00FF00", "italic": True},
                "code": {"font": "Mono"},
            }
        })
        assert style.run_overrides.italic.color == "#00FF00"
        assert style.run_overrides.italic.italic is True
        assert style.run_overrides.code.font == "Mono"
        # bold/link 用默认空 RunStyle
        assert style.run_overrides.bold.color is None


class TestParseRunOverrides:
    """测试 _parse_run_overrides。"""

    def test_empty_dict_returns_defaults(self) -> None:
        result = _parse_run_overrides({})
        assert result.bold.font is None
        assert result.italic.color is None

    def test_all_four_categories_parsed(self) -> None:
        data = {
            "bold": {"font": "Arial", "bold": True},
            "italic": {"font": "Times", "italic": True},
            "code": {"font": "Mono", "color": "#000000"},
            "link": {"color": "#0000FF", "underline": True},
        }
        result = _parse_run_overrides(data)
        assert result.bold.font == "Arial"
        assert result.bold.bold is True
        assert result.italic.font == "Times"
        assert result.code.font == "Mono"
        assert result.link.color == "#0000FF"
        assert result.link.underline is True


class TestParseRunStyle:
    """测试 _parse_run_style。"""

    def test_empty_dict_returns_all_none(self) -> None:
        result = _parse_run_style({})
        assert result.font is None
        assert result.font_size is None
        assert result.color is None
        assert result.bold is None
        assert result.italic is None
        assert result.underline is None
        assert result.background_color is None

    def test_all_fields_parsed(self) -> None:
        data = {
            "font": "Helvetica",
            "font_size": 14,
            "color": "#FF0000",
            "bold": True,
            "italic": False,
            "underline": True,
            "background_color": "#FFFF00",
        }
        result = _parse_run_style(data)
        assert result.font == "Helvetica"
        assert result.font_size == 14
        assert result.color == "#FF0000"
        assert result.bold is True
        assert result.italic is False
        assert result.underline is True
        assert result.background_color == "#FFFF00"


# ──────────────────────── layouts 配置解析测试 ────────────────────────


class TestLoadLayoutConfig:
    """测试 _load_layout_config。"""

    def test_loads_valid_layouts(self, valid_theme_pack: Path) -> None:
        manifest = _load_manifest(valid_theme_pack)
        config = _load_layout_config(valid_theme_pack, manifest)

        assert config.version == "1.0"
        assert len(config.layouts) == 2
        assert config.defaults.default == "title-and-content"

    def test_missing_layouts_raises_missing_file(self, tmp_path: Path) -> None:
        manifest = ThemePackManifest(name="t", version="1.0", author="a", files={})
        with pytest.raises(MissingFileError):
            _load_layout_config(tmp_path, manifest)


class TestParseLayoutConfig:
    """测试 _parse_layout_config。"""

    def test_minimal_config_uses_defaults(self) -> None:
        config = _parse_layout_config({})
        assert config.version == "1.0"
        assert isinstance(config.defaults, LayoutDefaults)
        assert config.defaults.default == "title-and-content"
        assert config.layouts == []
        assert config.groups == {}

    def test_groups_parsed_correctly(self) -> None:
        data = {
            "groups": {
                "content": {"display_name": "内容布局", "description": "内容类"},
                "section": {"display_name": "章节布局"},
            }
        }
        config = _parse_layout_config(data)
        assert "content" in config.groups
        assert config.groups["content"].display_name == "内容布局"
        assert config.groups["section"].display_name == "章节布局"

    def test_invalid_group_skipped(self) -> None:
        """非 dict 类型的 group 应被跳过。"""
        data = {"groups": {"bad": "not a dict"}}
        config = _parse_layout_config(data)
        assert "bad" not in config.groups

    def test_invalid_layout_skipped(self) -> None:
        """非 dict 类型的 layout 应被跳过。"""
        data = {"layouts": ["not a dict", 123]}
        config = _parse_layout_config(data)
        assert config.layouts == []


class TestParseLayoutDef:
    """测试 _parse_layout_def。"""

    def test_minimal_layout_def(self) -> None:
        data = {"id": "my-layout", "name": "My Layout"}
        layout = _parse_layout_def(data)
        assert layout.id == "my-layout"
        assert layout.name == "My Layout"
        assert layout.placeholders == []
        assert layout.keywords == []

    def test_full_layout_def_with_auto_apply_dict(self) -> None:
        data = {
            "id": "title-slide",
            "name": "Title Slide",
            "display_name": "标题页",
            "description": "标题幻灯片",
            "group": "basic",
            "placeholders": [
                {"index": 0, "type": "title", "role": "title"},
                {"index": 1, "type": "body", "name": "副标题"},
            ],
            "keywords": ["标题", "首页"],
            "tags": ["basic"],
            "auto_apply": {
                "conditions": ["first_slide"],
                "priority": 100,
            },
        }
        layout = _parse_layout_def(data)
        assert layout.id == "title-slide"
        assert layout.display_name == "标题页"
        assert len(layout.placeholders) == 2
        assert layout.placeholders[0].role == "title"
        assert layout.placeholders[1].name == "副标题"
        assert layout.keywords == ["标题", "首页"]
        assert layout.auto_apply is not None
        assert layout.auto_apply.conditions == ["first_slide"]
        assert layout.auto_apply.priority == 100

    def test_auto_apply_as_list_only_conditions(self) -> None:
        """auto_apply 为列表时应转换为 conditions。"""
        data = {
            "id": "x",
            "name": "X",
            "auto_apply": ["cond1", "cond2"],
        }
        layout = _parse_layout_def(data)
        assert layout.auto_apply is not None
        assert layout.auto_apply.conditions == ["cond1", "cond2"]
        assert layout.auto_apply.priority == 0

    def test_invalid_placeholder_skipped(self) -> None:
        """非 dict 类型的 placeholder 应被跳过。"""
        data = {
            "id": "x",
            "name": "X",
            "placeholders": ["not a dict", {"index": 0, "type": "title"}],
        }
        layout = _parse_layout_def(data)
        assert len(layout.placeholders) == 1


# ──────────────────────── 模板验证测试 ────────────────────────


class TestValidateTemplate:
    """测试 _validate_template。"""

    def test_valid_template_passes(self, valid_theme_pack: Path) -> None:
        manifest = _load_manifest(valid_theme_pack)
        layout_config = _load_layout_config(valid_theme_pack, manifest)
        # 不应抛异常
        _validate_template(valid_theme_pack / "template.pptx", layout_config)

    def test_missing_layout_in_template_raises(self, tmp_path: Path) -> None:
        """layouts.yaml 中定义了模板中不存在的布局应抛 InvalidConfigError。"""
        pptx_path = tmp_path / "template.pptx"
        _create_pptx_with_layouts(pptx_path, ["Title Slide"])

        from ppt_generator.core.models import LayoutDef
        layout_config = LayoutConfig(
            layouts=[
                LayoutDef(id="x", name="Title Slide"),
                LayoutDef(id="y", name="Nonexistent Layout"),
            ]
        )
        with pytest.raises(InvalidConfigError):
            _validate_template(pptx_path, layout_config)

    def test_corrupted_template_raises_template_load_error(self, tmp_path: Path) -> None:
        """损坏的模板文件应抛 TemplateLoadError。"""
        bad_pptx = tmp_path / "bad.pptx"
        bad_pptx.write_bytes(b"not a pptx file")
        with pytest.raises(TemplateLoadError):
            _validate_template(bad_pptx, LayoutConfig())


# ──────────────────────── list_available_themes 测试 ────────────────────────


class TestListAvailableThemes:
    """测试 list_available_themes。"""

    def test_nonexistent_dir_returns_empty(self, tmp_path: Path) -> None:
        result = list_available_themes(tmp_path / "nonexistent")
        assert result == []

    def test_empty_dir_returns_empty(self, tmp_path: Path) -> None:
        result = list_available_themes(tmp_path)
        assert result == []

    def test_lists_valid_theme_packs(self, valid_theme_pack: Path) -> None:
        result = list_available_themes(valid_theme_pack.parent)
        assert len(result) == 1
        assert result[0].name == "测试主题"

    def test_skips_invalid_theme_packs(self, tmp_path: Path) -> None:
        """无效的主题包子目录应被跳过，不抛异常。"""
        # 一个没有任何文件的子目录
        (tmp_path / "empty_pack").mkdir()
        # 一个 manifest 损坏的主题包
        bad_pack = tmp_path / "bad_pack"
        bad_pack.mkdir()
        (bad_pack / "manifest.yaml").write_text(
            "invalid: yaml: content:", encoding="utf-8"
        )
        # 一个合法的主题包
        valid_pack = tmp_path / "valid_pack"
        valid_pack.mkdir()
        (valid_pack / "manifest.yaml").write_text(
            "name: valid\nversion: 1.0\nauthor: a\n", encoding="utf-8"
        )

        result = list_available_themes(tmp_path)
        # 应只包含 valid_pack，跳过 bad_pack 和 empty_pack
        names = [m.name for m in result]
        assert "valid" in names
        assert len(result) == 1

    def test_ignores_files_in_themes_dir(self, tmp_path: Path) -> None:
        """themes_dir 中的普通文件应被忽略。"""
        (tmp_path / "readme.txt").write_text("not a theme", encoding="utf-8")
        result = list_available_themes(tmp_path)
        assert result == []

    def test_ignores_subdir_without_manifest(self, tmp_path: Path) -> None:
        """无 manifest.yaml 的子目录应被忽略。"""
        subdir = tmp_path / "no_manifest"
        subdir.mkdir()
        (subdir / "other.yaml").write_text("foo: bar", encoding="utf-8")
        result = list_available_themes(tmp_path)
        assert result == []
