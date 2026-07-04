"""集成测试。

本模块包含PPT生成完整管道的集成测试，验证端到端的生成流程。
"""

from pathlib import Path

import pytest
from returns.result import Failure, Success

from ppt_generator import PPTGenerator, MarkdownParseError, generate_ppt, parse_markdown, validate_slides
from ppt_generator.core.models import SlideSpec, SlideItem, SlideItemType


def test_generator_with_error_handling(tmp_path: Path) -> None:
    """测试生成器能够优雅地处理错误。"""
    from pptx import Presentation

    template_path = tmp_path / "template.pptx"
    Presentation().save(str(template_path))

    with pytest.raises(MarkdownParseError):
        generator = PPTGenerator(
            markdown_text="",
            template_path=template_path,
            output_path=tmp_path / "output.pptx",
        )
        generator.generate()


def test_generator_with_valid_content(tmp_path: Path) -> None:
    """测试完整生成流程与有效内容。"""
    from pptx import Presentation

    template_path = tmp_path / "template.pptx"
    output_path = tmp_path / "output.pptx"
    Presentation().save(str(template_path))

    markdown_text = "# 幻灯片1\n\n内容在这里\n\n# 幻灯片2\n\n更多内容"

    generator = PPTGenerator(
        markdown_text=markdown_text,
        template_path=template_path,
        output_path=output_path,
        title="测试演示文稿",
    )
    generator.generate()

    assert output_path.exists()
    assert generator.result is not None
    assert isinstance(generator.result, Success)


def test_model_validation() -> None:
    """测试模型验证输入是否正确。"""
    from ppt_generator import InvalidConfigError

    item = SlideItem(type=SlideItemType.PARAGRAPH, content="测试")
    assert item.type == SlideItemType.PARAGRAPH

    with pytest.raises(InvalidConfigError):
        SlideItem(type="invalid", content="测试")  # type: ignore

    spec = SlideSpec(title="测试", items=[item])
    assert spec.title == "测试"

    with pytest.raises(InvalidConfigError):
        SlideSpec(title="测试", items="不是列表")  # type: ignore


def test_functional_pipeline_parse_markdown() -> None:
    """测试函数式管道的parse_markdown函数。"""
    markdown_text = "# 幻灯片1\n\n内容"

    result = parse_markdown(markdown_text)

    assert isinstance(result, Success)
    slides = result.unwrap()
    assert len(slides) == 1
    assert slides[0].title == "幻灯片1"


def test_functional_pipeline_empty_markdown() -> None:
    """测试空Markdown解析。"""
    result = parse_markdown("")

    assert isinstance(result, Failure)
    assert isinstance(result.failure(), MarkdownParseError)


def test_functional_pipeline_validate_slides() -> None:
    """测试幻灯片验证函数。"""
    valid_slides = [SlideSpec(title="测试", items=[])]
    result = validate_slides(valid_slides)

    assert isinstance(result, Success)
    assert result.unwrap() == valid_slides


def test_functional_pipeline_validate_empty_slides() -> None:
    """测试空幻灯片列表验证。"""
    from ppt_generator import EmptySlideError

    result = validate_slides([])

    assert isinstance(result, Failure)
    assert isinstance(result.failure(), EmptySlideError)


def test_generate_ppt_function(tmp_path: Path) -> None:
    """测试generate_ppt函数式接口。"""
    from pptx import Presentation

    template_path = tmp_path / "template.pptx"
    output_path = tmp_path / "output.pptx"
    Presentation().save(str(template_path))

    markdown_text = "# 幻灯片1\n\n内容"

    result = generate_ppt(
        markdown_text=markdown_text,
        template_path=template_path,
        output_path=output_path,
        title="测试演示文稿",
    )

    assert isinstance(result, Success)
    assert result.unwrap() == output_path
    assert output_path.exists()
